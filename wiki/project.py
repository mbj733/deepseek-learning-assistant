# -*- coding: utf-8 -*-
"""项目管理（课程 = LLM Wiki 项目）。

每个课程在 sessions/{id}/ 下维护一个 wiki 项目。
"""

import os
import re
import shutil
import hashlib
from datetime import datetime
from pathlib import Path

from core.config import SESSIONS_DIR
from core.database import Database
from core.tokenizer import approx_tokens
from wiki import templates

PAGE_DIRS = ["sources", "entities", "topics", "synthesis"]


class Project:
    """单个课程项目的文件层管理。"""

    def __init__(self, session_id: str, name: str = ""):
        self.session_id = session_id
        self.root = Path(SESSIONS_DIR) / session_id
        self.raw_sources = self.root / "raw" / "sources"
        self.wiki_root = self.root / "wiki"
        self.db = Database()
        os.makedirs(self.raw_sources, exist_ok=True)
        os.makedirs(self.wiki_root, exist_ok=True)
        for d in PAGE_DIRS:
            os.makedirs(self.wiki_root / d, exist_ok=True)
        self.db.execute(
            "INSERT OR IGNORE INTO sessions (id, name, created_at, updated_at) VALUES (?,?,?,?)",
            (session_id, name or "未命名", datetime.now().isoformat(), datetime.now().isoformat()))
        self.db.commit()

    def init_wiki(self, topic: str = "课程", language: str = "zh"):
        """创建 schema/purpose/index/log/overview（幂等）。"""
        files = {
            "schema.md": templates.fill(templates.SCHEMA_TEMPLATE,
                                        topic=topic, language=language,
                                        wiki_root=str(self.wiki_root).replace("\\", "/")),
            "purpose.md": templates.fill(templates.PURPOSE_TEMPLATE, topic=topic),
            "index.md": templates.fill(templates.INDEX_TEMPLATE, sources="（暂无资料）"),
            "log.md": templates.LOG_TEMPLATE,
            "overview.md": templates.OVERVIEW_TEMPLATE,
        }
        for name, content in files.items():
            p = self.wiki_root / name
            if not p.exists():
                p.write_text(content, encoding="utf-8")
        self._append_log("初始化", "wiki", f"创建课程项目 {topic}")

    def ensure_wiki(self, topic: str = "课程", language: str = "zh"):
        if not (self.wiki_root / "schema.md").exists():
            self.init_wiki(topic, language)

    def import_file(self, src_path: str) -> tuple[bool, str, dict]:
        """复制资料进 raw/sources，写 FTS5 索引，返回文件信息。"""
        src = Path(src_path)
        if not src.exists():
            return False, "文件不存在", {}
        name = src.name
        dest = self.raw_sources / name
        base, ext = os.path.splitext(name)
        counter = 1
        while dest.exists():
            dest = self.raw_sources / f"{base}_{counter}{ext}"
            counter += 1
        shutil.copy2(str(src), str(dest))
        sha = sha256_file(str(dest))
        content = extract_text(str(dest))
        if not content or len(content.strip()) < 20:
            dest.unlink(missing_ok=True)
            return False, f"无法读取「{name}」的内容", {}
        chunks = chunk_text(content)
        now = datetime.now().isoformat()
        self.db.execute(
            "DELETE FROM knowledge_files WHERE session_id=? AND file_name=?",
            (self.session_id, name))
        self.db.execute(
            "DELETE FROM knowledge_chunks WHERE session_id=? AND source_file=?",
            (self.session_id, name))
        self.db.execute(
            """INSERT INTO knowledge_files
               (session_id, file_name, file_path, file_size, sha256, chunk_count, added_at)
               VALUES (?,?,?,?,?,?,?)""",
            (self.session_id, name, str(dest), dest.stat().st_size, sha, len(chunks), now))
        for i, c in enumerate(chunks):
            cur2 = self.db.execute(
                "INSERT INTO knowledge_chunks (session_id, source_file, chunk_index, content) VALUES (?,?,?,?)",
                (self.session_id, name, i, c))
            self.db.execute(
                "INSERT INTO knowledge_chunks_fts(rowid, content) VALUES (?,?)",
                (cur2.lastrowid, c))
        self.db.commit()
        info = {"name": name, "path": str(dest), "sha256": sha,
                "chunks": len(chunks), "tokens": approx_tokens(content)}
        return True, f"已导入「{name}」（{len(chunks)} 个片段）", info

    def remove_file(self, file_name: str) -> tuple[bool, str]:
        """删除 raw 文件 + 主表记录 + ingest 缓存。FTS 表只增不删，rebuild_fts() 清理。"""
        rows = self.db.fetchall(
            "SELECT file_path FROM knowledge_files WHERE session_id=? AND file_name=?",
            (self.session_id, file_name))
        for r in rows:
            p = r["file_path"]
            if p and os.path.exists(p):
                try:
                    os.remove(p)
                except OSError:
                    pass
        self.db.execute("DELETE FROM knowledge_files WHERE session_id=? AND file_name=?",
                        (self.session_id, file_name))
        self.db.execute("DELETE FROM knowledge_chunks WHERE session_id=? AND source_file=?",
                        (self.session_id, file_name))
        self.db.execute("DELETE FROM ingest_cache WHERE session_id=? AND source_path=?",
                        (self.session_id, file_name))
        self.db.commit()
        return True, f"已移除「{file_name}」"

    def rebuild_fts(self):
        """重建 FTS 索引（清理脏行）：DROP → CREATE → 全量回填。"""
        conn = self.db.get_conn()
        with conn:
            conn.execute("DROP TABLE IF EXISTS knowledge_chunks_fts")
            conn.execute("CREATE VIRTUAL TABLE knowledge_chunks_fts USING fts5(content, tokenize='unicode61')")
            rows = conn.execute(
                "SELECT id, content FROM knowledge_chunks WHERE session_id=?",
                (self.session_id,)).fetchall()
            for rid, content in rows:
                conn.execute("INSERT INTO knowledge_chunks_fts(rowid, content) VALUES (?,?)",
                             (rid, content))
        return len(rows)

    def list_wiki_pages(self) -> list[dict]:
        """返回 wiki 下所有 .md 页面（不含 schema/purpose 等控制文件）。"""
        pages = []
        for d in PAGE_DIRS:
            dpath = self.wiki_root / d
            if not dpath.exists():
                continue
            for f in sorted(dpath.glob("*.md")):
                pages.append({"rel": f"wiki/{d}/{f.name}", "type": d,
                              "title": f.stem, "path": str(f)})
        return pages

    def read_wiki_page(self, rel_path: str) -> str:
        fp = self.wiki_root / rel_path.replace("wiki/", "", 1)
        if fp.exists():
            return fp.read_text(encoding="utf-8")
        return ""

    def read_wiki(self, name: str) -> str:
        p = self.wiki_root / name
        return p.read_text(encoding="utf-8") if p.exists() else ""

    def write_wiki(self, name: str, content: str):
        (self.wiki_root / name).write_text(content, encoding="utf-8")

    def is_processed(self, file_name: str, sha: str) -> bool:
        row = self.db.fetchone(
            "SELECT sha256 FROM ingest_cache WHERE session_id=? AND source_path=?",
            (self.session_id, file_name))
        return bool(row) and row["sha256"] == sha

    def mark_processed(self, file_name: str, sha: str):
        self.db.execute(
            "INSERT OR REPLACE INTO ingest_cache (session_id, source_path, sha256, processed_at) VALUES (?,?,?,?)",
            (self.session_id, file_name, sha, datetime.now().isoformat()))
        self.db.commit()

    def pending_files(self) -> list[dict]:
        """raw/sources 下未处理或已变更的文件。"""
        result = []
        rows = self.db.fetchall(
            "SELECT file_name, file_path, sha256 FROM knowledge_files WHERE session_id=?",
            (self.session_id,))
        for r in rows:
            name = r["file_name"]
            sha = r["sha256"] or ""
            if not self.is_processed(name, sha):
                result.append({"name": name, "path": r["file_path"], "sha256": sha})
        return result

    def _append_log(self, action: str, obj: str, note: str = ""):
        line = f"| {datetime.now().strftime('%Y-%m-%d %H:%M')} | {action} | {obj} | {note} |\n"
        p = self.wiki_root / "log.md"
        content = p.read_text(encoding="utf-8") if p.exists() else templates.LOG_TEMPLATE
        p.write_text(content + line, encoding="utf-8")

    def append_log(self, action: str, obj: str, note: str = ""):
        self._append_log(action, obj, note)


def sha256_file(path: str) -> str:
    h = hashlib.sha256()
    with open(path, "rb") as f:
        for chunk in iter(lambda: f.read(65536), b""):
            h.update(chunk)
    return h.hexdigest()


def extract_text(filepath: str) -> str:
    """多格式文本提取：pdf/docx/pptx/txt/md。失败返回空串。"""
    ext = os.path.splitext(filepath)[1].lower()
    try:
        if ext == ".pdf":
            return _extract_pdf(filepath)
        if ext == ".docx":
            return _extract_docx(filepath)
        if ext == ".pptx":
            return _extract_pptx(filepath)
        if ext in (".txt", ".md", ".markdown", ".html", ".htm"):
            with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()
    except Exception:
        return ""
    return ""


def _extract_pdf(path: str) -> str:
    try:
        import fitz  # PyMuPDF
    except ImportError:
        return ""
    parts = []
    with fitz.open(path) as doc:
        for page in doc:
            parts.append(page.get_text())
    return "\n\n".join(parts)


def _extract_docx(path: str) -> str:
    try:
        import docx
    except ImportError:
        return ""
    d = docx.Document(path)
    parts = []
    for para in d.paragraphs:
        if para.text.strip():
            parts.append(para.text)
    for tbl in d.tables:
        for row in tbl.rows:
            parts.append(" | ".join(c.text.strip() for c in row.cells))
    return "\n".join(parts)


def _extract_pptx(path: str) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        return ""
    parts = []
    prs = Presentation(path)
    for slide in prs.slides:
        slide_parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_parts.append(shape.text)
        if slide_parts:
            parts.append("\n".join(slide_parts))
    return "\n\n---\n\n".join(parts)


CHUNK_SIZE = 600
CHUNK_OVERLAP = 150


def chunk_text(text: str) -> list[str]:
    """按段落优先、长度兜底的分块。"""
    text = re.sub(r"\r\n", "\n", text)
    paras = [p.strip() for p in re.split(r"\n\s*\n", text) if p.strip()]
    chunks: list[str] = []
    buf = ""
    for para in paras:
        for piece in _split_long_para(para):
            if len(buf) + len(piece) > CHUNK_SIZE and buf:
                chunks.append(buf)
                buf = piece
            else:
                buf = (buf + "\n" + piece) if buf else piece
    if buf:
        chunks.append(buf)
    return chunks or [text]


def _split_long_para(para: str) -> list[str]:
    if len(para) <= CHUNK_SIZE:
        return [para]
    pieces = []
    while len(para) > CHUNK_SIZE:
        cut = para[:CHUNK_SIZE]
        idx = max(cut.rfind("。"), cut.rfind(". "), cut.rfind("；"), cut.rfind(";"))
        if idx < CHUNK_SIZE // 2:
            idx = CHUNK_SIZE
        else:
            idx += 1
        pieces.append(para[:idx])
        para = para[idx - CHUNK_OVERLAP:]
    if para:
        pieces.append(para)
    return pieces
