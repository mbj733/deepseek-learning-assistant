#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
DeepSeek 学习助手 v2.0 - 多会话 + 知识库 RAG 系统

功能：
  - 多会话管理（左侧面板）
  - 每个会话有独立的记忆区（知识库），可上传 .txt/.pdf/.md
  - 问答严格基于上传的学习资料
  - DeepSeek API 流式输出
"""

import json
import os
import re
import threading
import tkinter as tk
from tkinter import ttk, messagebox, filedialog, scrolledtext
from datetime import datetime
from collections import Counter
from pathlib import Path
import shutil

import requests
import ttkbootstrap as ttkb
from ttkbootstrap.constants import *

# ── 路径 ──────────────────────────────────────────────────────────────
APP_DIR = os.path.dirname(os.path.abspath(__file__))
SESSIONS_DIR = os.path.join(APP_DIR, "sessions")
CONFIG_FILE = os.path.join(APP_DIR, "config.json")

DEEPSEEK_API_URL = "https://api.deepseek.com/chat/completions"
DEFAULT_MODEL = "deepseek-v4-flash"

# 确保目录存在
os.makedirs(SESSIONS_DIR, exist_ok=True)

# ── 系统提示词模板 ──────────────────────────────────────────────────
BASE_SYSTEM_PROMPT = """你是 DeepSeek 学习助手，一位专业的 AI 导师。

【核心规则 - 严格遵守】
1. 你的回答必须严格基于下面提供的【参考资料】，不得使用参考资料以外的知识
2. 如果用户的问题在参考资料中找不到答案，请明确回答"根据提供的资料，我无法回答这个问题"
3. 不要编造或捏造信息，不要使用你本身的预训练知识
4. 用清晰易懂的语言解释，辅以类比和例子（仅限参考资料中的内容）
5. 鼓励式教学，先肯定学生的思考再引导

【参考资料】
{context}

请基于以上参考资料回答用户的问题。如果参考资料为空，请告知用户先上传学习资料。"""

# ── 配置管理 ──────────────────────────────────────────────────────────
def load_config():
    if os.path.exists(CONFIG_FILE):
        try:
            with open(CONFIG_FILE, "r", encoding="utf-8") as f:
                return json.load(f)
        except:
            pass
    return {"api_key": "", "model": DEFAULT_MODEL}

def save_config(config):
    with open(CONFIG_FILE, "w", encoding="utf-8") as f:
        json.dump(config, f, ensure_ascii=False, indent=2)


# ── 知识库管理 ──────────────────────────────────────────────────────
CHUNK_SIZE = 500
CHUNK_OVERLAP = 100

class KnowledgeBase:
    """单个会话的知识库 - 管理上传的资料和检索"""

    def __init__(self, session_id):
        self.session_id = session_id
        self.data_dir = os.path.join(SESSIONS_DIR, session_id, "knowledge")
        self.meta_file = os.path.join(self.data_dir, "meta.json")
        self.files_dir = os.path.join(self.data_dir, "files")
        os.makedirs(self.files_dir, exist_ok=True)
        self._chunks = []
        self._file_meta = []
        self._load_meta()

    def _load_meta(self):
        """加载文件元数据"""
        if os.path.exists(self.meta_file):
            try:
                with open(self.meta_file, "r", encoding="utf-8") as f:
                    data = json.load(f)
                    self._file_meta = data.get("files", [])
            except:
                self._file_meta = []

    def _save_meta(self):
        """保存文件元数据"""
        with open(self.meta_file, "w", encoding="utf-8") as f:
            json.dump({"files": self._file_meta}, f, ensure_ascii=False, indent=2)

    def get_file_list(self):
        """获取已上传的文件列表"""
        return self._file_meta

    def add_file(self, src_path):
        """添加文件到知识库"""
        name = os.path.basename(src_path)
        # 避免重名
        base, ext = os.path.splitext(name)
        dest_path = os.path.join(self.files_dir, name)
        counter = 1
        while os.path.exists(dest_path):
            name = f"{base}_{counter}{ext}"
            dest_path = os.path.join(self.files_dir, name)
            counter += 1

        shutil.copy2(src_path, dest_path)

        # 提取文本
        content = self._extract_text(dest_path)
        if not content:
            os.remove(dest_path)
            return False, "无法读取文件内容"

        # 分块
        chunks = self._chunk_text(content, name)

        meta_entry = {
            "name": name,
            "path": dest_path,
            "size": os.path.getsize(dest_path),
            "chunks": len(chunks),
            "added_at": datetime.now().isoformat()
        }
        self._file_meta.append(meta_entry)
        self._chunks.extend(chunks)
        self._save_meta()
        return True, f"已添加 {name}（{len(chunks)} 个片段）"

    def remove_file(self, file_name):
        """删除文件"""
        for i, meta in enumerate(self._file_meta):
            if meta["name"] == file_name:
                if os.path.exists(meta["path"]):
                    os.remove(meta["path"])
                self._file_meta.pop(i)
                # 重建所有 chunks
                self._rebuild_chunks()
                self._save_meta()
                return True, f"已删除 {file_name}"
        return False, "文件不存在"

    def _rebuild_chunks(self):
        """重新构建所有块"""
        self._chunks = []
        for meta in self._file_meta:
            if os.path.exists(meta["path"]):
                content = self._extract_text(meta["path"])
                if content:
                    self._chunks.extend(self._chunk_text(content, meta["name"]))

    def _extract_text(self, filepath):
        """提取文件文本内容"""
        ext = os.path.splitext(filepath)[1].lower()

        if ext == ".txt":
            try:
                with open(filepath, "r", encoding="utf-8") as f:
                    return f.read()
            except:
                try:
                    with open(filepath, "r", encoding="gbk") as f:
                        return f.read()
                except:
                    return ""

        elif ext == ".md":
            try:
                with open(filepath, "r", encoding="utf-8") as f:
                    return f.read()
            except:
                return ""

        elif ext == ".pdf":
            return self._extract_pdf(filepath)

        elif ext == ".docx":
            return self._extract_docx(filepath)

        elif ext == ".pptx":
            return self._extract_pptx(filepath)

        return ""

    def _extract_pdf(self, filepath):
        """提取 PDF 文本（使用基本方法）"""
        try:
            import fitz  # PyMuPDF
            text_parts = []
            doc = fitz.open(filepath)
            for page in doc:
                text_parts.append(page.get_text())
            doc.close()
            return "\n".join(text_parts)
        except ImportError:
            # 尝试用 pdfminer
            try:
                from pdfminer.high_level import extract_text as pdf_extract
                return pdf_extract(filepath)
            except ImportError:
                return ""
        except Exception:
            return ""

    def _extract_docx(self, filepath):
        """提取 Word 文档文本"""
        try:
            from docx import Document
            doc = Document(filepath)
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            # 也提取表格中的文字
            tables_text = []
            for table in doc.tables:
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        tables_text.append(" | ".join(cells))
            all_text = "\n".join(paragraphs)
            if tables_text:
                all_text += "\n\n[表格内容]\n" + "\n".join(tables_text)
            return all_text
        except Exception:
            return ""

    def _extract_pptx(self, filepath):
        """提取 PowerPoint 文本"""
        try:
            from pptx import Presentation
            prs = Presentation(filepath)
            text_parts = []
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            if para.text.strip():
                                slide_texts.append(para.text.strip())
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                            if cells:
                                slide_texts.append(" | ".join(cells))
                if slide_texts:
                    text_parts.append(f"[第{slide_num}页]\n" + "\n".join(slide_texts))
            return "\n\n".join(text_parts)
        except Exception:
            return ""

    def _chunk_text(self, text, source_name):
        """将文本分割为块"""
        if not text:
            return []

        # 按段落分割
        paragraphs = re.split(r'\n\s*\n', text)
        chunks = []
        current = ""

        for para in paragraphs:
            para = para.strip()
            if not para:
                continue

            if len(current) + len(para) < CHUNK_SIZE:
                if current:
                    current += "\n\n" + para
                else:
                    current = para
            else:
                if current:
                    chunks.append({
                        "text": current,
                        "source": source_name
                    })
                # 如果段落本身就很长，再细分
                if len(para) > CHUNK_SIZE:
                    for i in range(0, len(para), CHUNK_SIZE - CHUNK_OVERLAP):
                        chunk_text = para[i:i + CHUNK_SIZE]
                        if len(chunk_text) > 50:  # 忽略太短的碎片
                            chunks.append({
                                "text": chunk_text,
                                "source": source_name
                            })
                else:
                    current = para

        if current:
            chunks.append({
                "text": current,
                "source": source_name
            })

        return chunks

    def search(self, query, top_k=5):
        """搜索最相关的知识块（基于关键词匹配 + 排序）"""
        if not self._chunks:
            return []

        query_lower = query.lower()
        query_words = set(re.findall(r'\w+', query_lower))

        if not query_words:
            return []

        scored = []
        for chunk in self._chunks:
            text_lower = chunk["text"].lower()
            text_words = set(re.findall(r'\w+', text_lower))

            # 计算交集
            common = query_words & text_words
            if not common:
                continue

            # 得分 = 匹配词数量 + 稀有词加权
            score = len(common)

            # 精确匹配加分
            for q_word in query_words:
                if len(q_word) > 2:
                    count = text_lower.count(q_word)
                    score += min(count, 5) * 0.5

            # 短语匹配加分（2-4 个连续词）
            query_phrases = re.findall(r'(?:\S+\s){1,3}\S+', query_lower)
            for phrase in query_phrases:
                if phrase in text_lower:
                    score += 3

            scored.append((score, chunk))

        # 排序取 top_k
        scored.sort(key=lambda x: -x[0])
        top_chunks = scored[:top_k]

        # 去重（相同文本只保留一次）
        seen_texts = set()
        deduped = []
        for _, chunk in top_chunks:
            if chunk["text"][:100] not in seen_texts:
                seen_texts.add(chunk["text"][:100])
                deduped.append(chunk)

        return deduped

    def format_context(self, query):
        """将搜索结果格式化为上下文文本"""
        results = self.search(query)
        if not results:
            return "（当前会话未上传学习资料，或资料中未找到相关内容）"

        parts = []
        for i, chunk in enumerate(results, 1):
            source = chunk.get("source", "未知来源")
            parts.append(f"[片段 {i} - 来自 {source}]\n{chunk['text']}\n")

        return "\n---\n".join(parts)

    def get_context_stat(self):
        """获取知识库统计"""
        file_count = len(self._file_meta)
        chunk_count = len(self._chunks)
        total_chars = sum(len(c["text"]) for c in self._chunks)
        return file_count, chunk_count, total_chars


# ── 会话管理 ──────────────────────────────────────────────────────────
class SessionManager:
    """多会话管理器"""

    def __init__(self):
        self.sessions_file = os.path.join(SESSIONS_DIR, "_sessions.json")
        self.sessions = {}
        self._load()

    def _load(self):
        """加载会话列表"""
        if os.path.exists(self.sessions_file):
            try:
                with open(self.sessions_file, "r", encoding="utf-8") as f:
                    self.sessions = json.load(f)
            except:
                self.sessions = {}
        if not self.sessions:
            # 创建默认会话
            default_id = self._generate_id()
            self.sessions[default_id] = {
                "name": "默认课程",
                "created_at": datetime.now().isoformat(),
                "updated_at": datetime.now().isoformat(),
                "message_count": 0
            }
            self._save()

    def _save(self):
        """保存会话列表"""
        os.makedirs(SESSIONS_DIR, exist_ok=True)
        with open(self.sessions_file, "w", encoding="utf-8") as f:
            json.dump(self.sessions, f, ensure_ascii=False, indent=2)

    def _generate_id(self):
        return datetime.now().strftime("%Y%m%d%H%M%S") + str(os.urandom(2).hex())

    def get_all(self):
        """获取所有会话"""
        return self.sessions

    def create(self, name="新课程"):
        """创建新会话"""
        sid = self._generate_id()
        self.sessions[sid] = {
            "name": name,
            "created_at": datetime.now().isoformat(),
            "updated_at": datetime.now().isoformat(),
            "message_count": 0
        }
        # 创建知识库目录
        kb_dir = os.path.join(SESSIONS_DIR, sid, "knowledge")
        os.makedirs(kb_dir, exist_ok=True)
        self._save()
        return sid

    def rename(self, sid, new_name):
        """重命名会话"""
        if sid in self.sessions:
            self.sessions[sid]["name"] = new_name
            self._save()
            return True
        return False

    def delete(self, sid):
        """删除会话及其知识库"""
        if sid in self.sessions:
            # 删除文件
            session_dir = os.path.join(SESSIONS_DIR, sid)
            if os.path.exists(session_dir):
                shutil.rmtree(session_dir)
            del self.sessions[sid]
            self._save()
            return True
        return False

    def update_time(self, sid):
        """更新最后使用时间"""
        if sid in self.sessions:
            self.sessions[sid]["updated_at"] = datetime.now().isoformat()
            self.sessions[sid]["message_count"] = self.sessions[sid].get("message_count", 0) + 1
            self._save()

    def get_name(self, sid):
        return self.sessions.get(sid, {}).get("name", "未命名")


# ── DeepSeek API ─────────────────────────────────────────────────────
class DeepSeekClient:
    def __init__(self, api_key, model=DEFAULT_MODEL):
        self.api_key = api_key
        self.model = model
        self.headers = {
            "Authorization": f"Bearer {api_key}",
            "Content-Type": "application/json"
        }

    def chat_stream(self, messages):
        """流式聊天"""
        payload = {
            "model": self.model,
            "messages": messages,
            "stream": True,
            "temperature": 0.3,  # 降低温度，让回答更忠于资料
            "max_tokens": 4096,
        }
        response = requests.post(
            DEEPSEEK_API_URL,
            headers=self.headers,
            json=payload,
            stream=True,
            timeout=120
        )
        response.raise_for_status()

        for line in response.iter_lines(decode_unicode=True):
            if line:
                if line.startswith("data: "):
                    data_str = line[6:]
                    if data_str.strip() == "[DONE]":
                        break
                    try:
                        data = json.loads(data_str)
                        delta = data["choices"][0].get("delta", {})
                        content = delta.get("content", "")
                        if content:
                            yield content
                    except (json.JSONDecodeError, KeyError):
                        continue


# ══════════════════════════════════════════════════════════════════════
#  GUI 主应用
# ══════════════════════════════════════════════════════════════════════

class DeepSeekLearnerApp:
    def __init__(self):
        self.config = load_config()
        self.client = None
        self.session_manager = SessionManager()
        self.current_sid = None
        self.kb = None
        self.messages = []  # 当前会话的消息历史
        self.is_loading = False
        self.current_ai_widget = None
        self.current_ai_container = None

        # 选择第一个会话
        sessions = self.session_manager.get_all()
        self.current_sid = next(iter(sessions.keys()))
        self.kb = KnowledgeBase(self.current_sid)

        # 窗口
        self.root = ttkb.Window(
            title="📚 DeepSeek 学习助手 v2.0",
            themename="litera",
            size=(1100, 720),
            minsize=(800, 500),
        )

        self._setup_ui()
        self._update_session_display()

        if self.config.get("api_key"):
            self.client = DeepSeekClient(
                self.config["api_key"],
                self.config.get("model", DEFAULT_MODEL)
            )
            self.status_label.configure(text="● 已连接", bootstyle="success")
            self._refresh_kb_display()
        else:
            self._add_system_message("👋 欢迎！请先点击 ⚙ 设置 API Key，然后上传学习资料开始学习")

        self.root.protocol("WM_DELETE_WINDOW", self._on_close)

    # ════════════════════════════════════════════════════════════════
    #  UI 构建
    # ════════════════════════════════════════════════════════════════

    def _setup_ui(self):
        """构建主界面"""
        # 主布局：左侧栏 + 右侧
        main_pw = ttk.PanedWindow(self.root, orient=HORIZONTAL)
        main_pw.pack(fill=BOTH, expand=True)

        # ── 左侧面板 ──
        left_frame = ttkb.Frame(main_pw, width=300)
        left_frame.pack_propagate(False)

        # 左侧标题
        title_bar = ttkb.Frame(left_frame, padding=(10, 8))
        title_bar.pack(fill=X)
        ttkb.Label(
            title_bar, text="📚 课程列表",
            font=("Microsoft YaHei UI", 12, "bold")
        ).pack(side=LEFT)
        ttkb.Button(
            title_bar, text="＋",
            bootstyle="success-outline",
            width=3,
            command=self._create_session
        ).pack(side=RIGHT)

        # 会话列表（使用 Canvas + Scrollbar 实现滚动）
        session_list_frame = ttkb.Frame(left_frame)
        session_list_frame.pack(fill=BOTH, expand=True, padx=5, pady=5)

        self.session_canvas = tk.Canvas(session_list_frame, highlightthickness=0, bg="#f8f9fa")
        session_scroll = ttkb.Scrollbar(session_list_frame, orient=VERTICAL, command=self.session_canvas.yview)
        self.session_canvas.configure(yscrollcommand=session_scroll.set)

        self.session_list_inner = ttkb.Frame(self.session_canvas)
        self.session_list_inner.bind("<Configure>",
            lambda e: self.session_canvas.configure(scrollregion=self.session_canvas.bbox("all")))

        self.session_canvas_window = self.session_canvas.create_window(
            (0, 0), window=self.session_list_inner, anchor="nw", width=270)

        self.session_canvas.pack(side=LEFT, fill=BOTH, expand=True)
        session_scroll.pack(side=RIGHT, fill=Y)

        # 分隔线
        ttk.Separator(left_frame, orient=HORIZONTAL).pack(fill=X, padx=10, pady=5)

        # 当前会话的知识库信息
        kb_title_frame = ttkb.Frame(left_frame)
        kb_title_frame.pack(fill=X, padx=10, pady=(5, 0))
        ttkb.Label(
            kb_title_frame, text="📄 知识库",
            font=("Microsoft YaHei UI", 10, "bold")
        ).pack(side=LEFT)

        # 上传按钮
        upload_btn = ttkb.Button(
            kb_title_frame, text="📤 上传文件",
            bootstyle="info-outline",
            width=10,
            command=self._upload_file
        )
        upload_btn.pack(side=RIGHT)

        # 知识库文件列表
        kb_list_frame = ttkb.Frame(left_frame)
        kb_list_frame.pack(fill=BOTH, expand=True, padx=10, pady=5)

        self.kb_canvas = tk.Canvas(kb_list_frame, height=200, highlightthickness=0, bg="#f8f9fa")
        kb_scroll = ttkb.Scrollbar(kb_list_frame, orient=VERTICAL, command=self.kb_canvas.yview)
        self.kb_canvas.configure(yscrollcommand=kb_scroll.set)

        self.kb_list_inner = ttkb.Frame(self.kb_canvas)
        self.kb_list_inner.bind("<Configure>",
            lambda e: self.kb_canvas.configure(scrollregion=self.kb_canvas.bbox("all")))

        self.kb_canvas_window = self.kb_canvas.create_window(
            (0, 0), window=self.kb_list_inner, anchor="nw", width=265)

        self.kb_canvas.pack(side=LEFT, fill=BOTH, expand=True)
        kb_scroll.pack(side=RIGHT, fill=Y)

        # 知识库状态
        self.kb_status_label = ttkb.Label(
            left_frame, text="0 个文件",
            font=("Microsoft YaHei UI", 9),
            bootstyle="secondary"
        )
        self.kb_status_label.pack(fill=X, padx=10, pady=(0, 8))

        left_frame.add_to_pw = lambda: main_pw.insert(0, left_frame)
        main_pw.add(left_frame, weight=0)

        # ── 右侧主区域 ──
        right_frame = ttkb.Frame(main_pw)
        right_frame.grid_rowconfigure(1, weight=1)
        right_frame.grid_columnconfigure(0, weight=1)

        # 标题栏
        title_bar = ttkb.Frame(right_frame, padding=(15, 6))
        title_bar.grid(row=0, column=0, sticky="ew")
        title_bar.grid_columnconfigure(1, weight=1)

        self.session_title_label = ttkb.Label(
            title_bar, text="📖 默认课程",
            font=("Microsoft YaHei UI", 14, "bold")
        )
        self.session_title_label.pack(side=LEFT)

        # 状态
        self.status_label = ttkb.Label(
            title_bar, text="● 未连接",
            font=("Microsoft YaHei UI", 9),
            bootstyle="secondary"
        )
        self.status_label.pack(side=RIGHT, padx=10)

        # 设置按钮
        ttkb.Button(
            title_bar, text="⚙ 设置",
            bootstyle="outline-secondary",
            command=self.open_settings,
            width=7
        ).pack(side=RIGHT, padx=3)

        # 聊天区
        chat_frame = ttkb.Frame(right_frame, padding=(10, 0))
        chat_frame.grid(row=1, column=0, sticky="nsew")
        chat_frame.grid_rowconfigure(0, weight=1)
        chat_frame.grid_columnconfigure(0, weight=1)

        self.chat_canvas = tk.Canvas(chat_frame, highlightthickness=0, bg="#f0f2f5")
        chat_scroll = ttk.Scrollbar(chat_frame, orient=VERTICAL, command=self.chat_canvas.yview)
        self.chat_canvas.configure(yscrollcommand=chat_scroll.set)

        self.msg_frame = ttkb.Frame(self.chat_canvas, padding=10)
        self.msg_frame.bind("<Configure>",
            lambda e: self.chat_canvas.configure(scrollregion=self.chat_canvas.bbox("all")))

        self.canvas_window = self.chat_canvas.create_window(
            (0, 0), window=self.msg_frame, anchor="nw",
            width=self.chat_canvas.winfo_width())
        self.chat_canvas.bind("<Configure>", self._on_canvas_resize)

        self.chat_canvas.grid(row=0, column=0, sticky="nsew")
        chat_scroll.grid(row=0, column=1, sticky="ns")

        # 输入区
        input_frame = ttkb.Frame(right_frame, padding=(10, 8))
        input_frame.grid(row=2, column=0, sticky="ew")
        input_frame.grid_columnconfigure(0, weight=1)

        self.input_text = scrolledtext.ScrolledText(
            input_frame,
            height=3,
            wrap=WORD,
            font=("Microsoft YaHei UI", 10),
            padx=10, pady=8,
            relief="flat",
            borderwidth=1,
        )
        self.input_text.grid(row=0, column=0, sticky="ew", padx=(0, 8))
        self.input_text.bind("<Return>", self._on_enter_key)

        self.send_btn = ttkb.Button(
            input_frame, text="发送 ▶",
            bootstyle="success",
            command=self.send_message,
            width=10,
        )
        self.send_btn.grid(row=0, column=1, sticky="ns")

        ttkb.Label(
            input_frame,
            text="Enter 发送 | Shift+Enter 换行",
            font=("Microsoft YaHei UI", 8),
            bootstyle="secondary"
        ).grid(row=1, column=0, columnspan=2, sticky="w", pady=(2, 0))

        main_pw.add(right_frame, weight=1)

    def _on_canvas_resize(self, event):
        self.chat_canvas.itemconfig(self.canvas_window, width=event.width)

    # ════════════════════════════════════════════════════════════════
    #  会话列表
    # ════════════════════════════════════════════════════════════════

    def _update_session_display(self):
        """刷新左侧会话列表"""
        for w in self.session_list_inner.winfo_children():
            w.destroy()

        sessions = self.session_manager.get_all()
        # 按更新时间倒序
        sorted_items = sorted(sessions.items(), key=lambda x: x[1].get("updated_at", ""), reverse=True)

        for sid, info in sorted_items:
            is_active = sid == self.current_sid
            frame = ttkb.Frame(self.session_list_inner, padding=(5, 4))
            frame.pack(fill=X, pady=1)

            # 选中高亮
            if is_active:
                frame.configure(bootstyle="primary")

            # 图标 + 名称
            name = info.get("name", "未命名")
            count = info.get("message_count", 0)

            name_label = ttkb.Label(
                frame, text=f"📁 {name}",
                font=("Microsoft YaHei UI", 10, "bold" if is_active else "normal"),
                bootstyle="inverse-primary" if is_active else ""
            )
            name_label.pack(side=LEFT, fill=X, expand=True)

            # 右键菜单绑定
            for widget in [frame, name_label]:
                widget.bind("<Button-3>", lambda e, s=sid, n=name: self._show_session_menu(e, s, n))
                widget.bind("<Button-1>", lambda e, s=sid: self._switch_session(s))

            # 消息数
            ttkb.Label(
                frame, text=str(count),
                font=("Microsoft YaHei UI", 9),
                bootstyle="secondary"
            ).pack(side=RIGHT, padx=(2, 0))

    def _show_session_menu(self, event, sid, name):
        """右键菜单 - 会话操作"""
        menu = tk.Menu(self.root, tearoff=0, font=("Microsoft YaHei UI", 9))
        menu.add_command(label="✏️ 重命名", command=lambda: self._rename_session(sid))
        menu.add_command(label="🗑 删除", command=lambda: self._delete_session(sid))
        menu.tk_popup(event.x_root, event.y_root)

    def _create_session(self, name=None):
        """新建会话"""
        if not name:
            name = f"课程 {len(self.session_manager.get_all()) + 1}"
        sid = self.session_manager.create(name)
        self.current_sid = sid
        self.kb = KnowledgeBase(sid)
        self.messages = []
        self._update_session_display()
        self._refresh_chat_area()
        self._update_session_title()
        self._refresh_kb_display()
        self._add_system_message(f"📖 已创建新课程「{name}」，请上传学习资料开始学习")

    def _rename_session(self, sid):
        """重命名会话"""
        dialog = ttkb.Toplevel(self.root)
        dialog.title("重命名")
        dialog.geometry("350x130")
        dialog.resizable(False, False)
        dialog.transient(self.root)
        dialog.grab_set()

        ttkb.Label(dialog, text="新名称：", font=("Microsoft YaHei UI", 10)).pack(padx=20, pady=(20, 5))
        entry = ttkb.Entry(dialog, font=("Microsoft YaHei UI", 10))
        entry.pack(fill=X, padx=20, pady=5)
        entry.insert(0, self.session_manager.get_name(sid))
        entry.focus()
        entry.select_range(0, END)

        def do_rename():
            new_name = entry.get().strip()
            if new_name:
                self.session_manager.rename(sid, new_name)
                self._update_session_display()
                self._update_session_title()
            dialog.destroy()

        ttkb.Button(dialog, text="确定", bootstyle="success", command=do_rename).pack(pady=10)

        dialog.bind("<Return>", lambda e: do_rename())

        # 居中
        dialog.update_idletasks()
        x = self.root.winfo_x() + (self.root.winfo_width() - dialog.winfo_width()) // 2
        y = self.root.winfo_y() + (self.root.winfo_height() - dialog.winfo_height()) // 2
        dialog.geometry(f"+{x}+{y}")

    def _delete_session(self, sid):
        """删除会话"""
        sessions = self.session_manager.get_all()
        if len(sessions) <= 1:
            messagebox.showwarning("提示", "至少保留一个会话")
            return
        if not messagebox.askyesno("确认删除", f"确定要删除「{self.session_manager.get_name(sid)}」吗？\n（知识库文件也会被删除）"):
            return

        self.session_manager.delete(sid)
        # 切换到另一个会话
        remaining = self.session_manager.get_all()
        self.current_sid = next(iter(remaining.keys()))
        self.kb = KnowledgeBase(self.current_sid)
        self.messages = []
        self._update_session_display()
        self._refresh_chat_area()
        self._update_session_title()
        self._refresh_kb_display()
        self._add_system_message(f"已切换到「{self.session_manager.get_name(self.current_sid)}」")

    def _switch_session(self, sid):
        """切换会话"""
        if sid == self.current_sid or self.is_loading:
            return
        self.current_sid = sid
        self.kb = KnowledgeBase(sid)
        self.messages = []
        self._update_session_display()
        self._refresh_chat_area()
        self._update_session_title()
        self._refresh_kb_display()
        self._add_system_message(f"已切换到「{self.session_manager.get_name(sid)}」")

    def _update_session_title(self):
        """更新当前会话标题"""
        name = self.session_manager.get_name(self.current_sid)
        self.session_title_label.configure(text=f"📖 {name}")

    # ════════════════════════════════════════════════════════════════
    #  知识库显示
    # ════════════════════════════════════════════════════════════════

    def _refresh_kb_display(self):
        """刷新知识库文件列表"""
        # 清空
        for w in self.kb_list_inner.winfo_children():
            w.destroy()

        files = self.kb.get_file_list()

        if not files:
            ttkb.Label(
                self.kb_list_inner,
                text="📭 暂无资料\n点击「上传」添加学习资料",
                font=("Microsoft YaHei UI", 9),
                bootstyle="secondary",
                justify=CENTER
            ).pack(pady=20)
        else:
            for meta in files:
                file_frame = ttkb.Frame(self.kb_list_inner)
                file_frame.pack(fill=X, pady=1)

                icon = "📄"
                if meta["name"].lower().endswith(".pdf"):
                    icon = "📕"
                elif meta["name"].lower().endswith(".txt"):
                    icon = "📝"
                elif meta["name"].lower().endswith(".md"):
                    icon = "📋"
                elif meta["name"].lower().endswith(".docx"):
                    icon = "📘"
                elif meta["name"].lower().endswith(".pptx"):
                    icon = "📊"

                ttkb.Label(
                    file_frame, text=f"{icon} {meta['name']}  ({meta['chunks']}段)",
                    font=("Microsoft YaHei UI", 9),
                ).pack(side=LEFT, fill=X, expand=True)

                # 删除按钮
                ttkb.Button(
                    file_frame, text="✕",
                    bootstyle="danger-outline",
                    width=2,
                    command=lambda n=meta['name']: self._remove_file(n)
                ).pack(side=RIGHT)

        # 更新统计
        file_count, chunk_count, total_chars = self.kb.get_context_stat()
        kb_size = ""
        if total_chars > 10000:
            kb_size = f"（{total_chars // 10000}万字）"
        self.kb_status_label.configure(
            text=f"{file_count} 个文件 · {chunk_count} 个片段{kb_size}"
        )

    def _upload_file(self):
        """上传文件到当前会话的知识库"""
        filetypes = [
            ("支持的文件", "*.txt *.md *.pdf *.docx *.pptx"),
            ("文本文件", "*.txt"),
            ("Markdown 文档", "*.md"),
            ("PDF 文档", "*.pdf"),
            ("Word 文档", "*.docx"),
            ("PowerPoint 演示", "*.pptx"),
            ("所有文件", "*.*"),
        ]
        files = filedialog.askopenfilenames(
            title="选择学习资料",
            filetypes=filetypes
        )
        if not files:
            return

        for f in files:
            success, msg = self.kb.add_file(f)
            if not success:
                messagebox.showwarning("上传失败", msg)

        self._refresh_kb_display()
        kb_name = os.path.basename(files[0]) if len(files) == 1 else f"{len(files)} 个文件"
        self._add_system_message(f"📎 已上传 {kb_name}，现在可以基于资料提问了")

    def _remove_file(self, file_name):
        """删除知识库中的文件"""
        if messagebox.askyesno("确认删除", f"从知识库中移除「{file_name}」？"):
            self.kb.remove_file(file_name)
            self._refresh_kb_display()
            self._add_system_message(f"🗑 已移除「{file_name}」")

    # ════════════════════════════════════════════════════════════════
    #  聊天区
    # ════════════════════════════════════════════════════════════════

    def _refresh_chat_area(self):
        """清空聊天区"""
        for w in self.msg_frame.winfo_children():
            w.destroy()

    def _add_system_message(self, text):
        """添加系统消息"""
        frame = ttkb.Frame(self.msg_frame)
        frame.pack(fill=X, pady=4, padx=5)

        inner = ttkb.Frame(frame)
        inner.pack(anchor="center")

        ttkb.Label(
            inner, text=f"📌 {text}",
            font=("Microsoft YaHei UI", 9),
            bootstyle="secondary",
        ).pack(padx=10, pady=5)

    def _create_bubble(self, role, content, is_streaming=False):
        """创建聊天气泡"""
        is_user = role == "user"

        frame = ttkb.Frame(self.msg_frame)
        frame.pack(fill=X, pady=(0, 8), padx=5)

        inner = ttkb.Frame(frame)
        if is_user:
            inner.pack(anchor="e")
        else:
            inner.pack(anchor="w")

        # 标签
        tag_text = "👤 你" if is_user else "🤖 DeepSeek (基于资料)"
        tag_color = "#4caf50" if is_user else "#2196f3"
        ttkb.Label(
            inner, text=tag_text,
            font=("Microsoft YaHei UI", 9, "bold"),
            foreground=tag_color,
        ).pack(anchor="w" if not is_user else "e", padx=5, pady=(0, 2))

        # 气泡
        bg = "#d4f0d4" if is_user else "#ffffff"

        text_w = tk.Text(
            inner,
            wrap=WORD,
            font=("Microsoft YaHei UI", 10),
            bg=bg,
            relief="flat",
            padx=14, pady=10,
            height=3, width=60,
            highlightthickness=0,
            borderwidth=1,
        )
        text_w.pack()
        text_w.insert("1.0", content)
        if not is_streaming:
            text_w.configure(state=DISABLED)
        else:
            text_w.configure(state="normal")

        # 时间戳
        time_str = datetime.now().strftime("%H:%M")
        anchor = "e" if is_user else "w"
        ttkb.Label(
            inner, text=time_str,
            font=("Microsoft YaHei UI", 8),
            bootstyle="secondary"
        ).pack(anchor=anchor, padx=5, pady=(2, 0))

        return text_w

    # ════════════════════════════════════════════════════════════════
    #  发送消息
    # ════════════════════════════════════════════════════════════════

    def send_message(self):
        if self.is_loading:
            return

        text = self.input_text.get("1.0", "end-1c").strip()
        if not text:
            return

        if not self.client or not self.config.get("api_key"):
            messagebox.showwarning("提示", "请先设置 DeepSeek API Key")
            self.open_settings()
            return

        # 清空输入
        self.input_text.delete("1.0", END)

        # 显示用户消息
        self._create_bubble("user", text)
        self.messages.append({"role": "user", "content": text})

        # 滚动
        self.root.after(50, self._scroll_to_bottom)

        # 调用 API
        self.is_loading = True
        self.send_btn.configure(text="响应中...", state=DISABLED)
        self._update_status("思考中...", "warning")

        thread = threading.Thread(target=self._do_chat, args=(text,), daemon=True)
        thread.start()

    def _do_chat(self, user_text):
        """在线程中执行 API 调用（RAG 模式）"""
        try:
            # 1. 从知识库检索相关内容
            context = self.kb.format_context(user_text)

            # 2. 构建系统提示词（基于资料）
            system_prompt = BASE_SYSTEM_PROMPT.format(context=context)

            # 3. 构建消息列表
            messages = [{"role": "system", "content": system_prompt}]
            messages.extend(self.messages)

            # 4. 创建流式气泡
            self.root.after(0, self._create_streaming_bubble)

            # 5. 流式请求
            full_response = ""
            for chunk in self.client.chat_stream(messages):
                full_response += chunk
                self.root.after(0, self._update_streaming_text, full_response)

            # 6. 完成
            self.root.after(0, self._finish_chat, full_response)
            self.session_manager.update_time(self.current_sid)
            self._update_session_display()

        except requests.exceptions.HTTPError as e:
            msg = f"API 请求失败 (HTTP {e.response.status_code})"
            if e.response.status_code == 401:
                msg = "API Key 无效，请检查设置"
            elif e.response.status_code == 429:
                msg = "请求过于频繁，请稍后再试"
            self.root.after(0, self._handle_error, msg)

        except requests.exceptions.ConnectionError:
            self.root.after(0, self._handle_error, "无法连接到 DeepSeek API，请检查网络")

        except requests.exceptions.Timeout:
            self.root.after(0, self._handle_error, "请求超时，请稍后重试")

        except Exception as e:
            self.root.after(0, self._handle_error, f"发生错误：{str(e)}")

    def _create_streaming_bubble(self):
        """创建流式输出气泡"""
        inner = ttkb.Frame(self.msg_frame)
        inner.pack(fill=X, pady=(0, 8), padx=5)

        container = ttkb.Frame(inner)
        container.pack(anchor="w")

        ttkb.Label(
            container, text="🤖 DeepSeek (基于资料)",
            font=("Microsoft YaHei UI", 9, "bold"),
            foreground="#2196f3",
        ).pack(anchor="w", padx=5, pady=(0, 2))

        text_w = tk.Text(
            container,
            wrap=WORD,
            font=("Microsoft YaHei UI", 10),
            bg="#ffffff",
            relief="flat",
            padx=14, pady=10,
            height=3, width=60,
            highlightthickness=0,
            borderwidth=1,
        )
        text_w.pack()
        text_w.configure(state="normal")
        text_w.insert("1.0", "▊")

        self.current_ai_widget = text_w
        self.current_ai_container = container

        self.root.after(300, self._blink_cursor)
        self._scroll_to_bottom()

    def _blink_cursor(self):
        """闪烁光标"""
        if self.current_ai_widget and self.is_loading:
            content = self.current_ai_widget.get("1.0", "end-1c")
            if content.endswith("▊"):
                self.current_ai_widget.delete("end-2c", "end-1c")
                self.current_ai_widget.insert(END, "▌")
            elif content.endswith("▌"):
                self.current_ai_widget.delete("end-2c", "end-1c")
                self.current_ai_widget.insert(END, "▊")
            self.root.after(400, self._blink_cursor)

    def _update_streaming_text(self, text):
        """更新流式文本"""
        if self.current_ai_widget:
            self.current_ai_widget.configure(state="normal")
            self.current_ai_widget.delete("1.0", END)
            self.current_ai_widget.insert("1.0", text + "▊")

            # 自适应高度
            lines = int(self.current_ai_widget.index("end-1c").split(".")[0])
            self.current_ai_widget.configure(height=min(max(lines + 1, 3), 22))

            self._scroll_to_bottom()

    def _finish_chat(self, full_response):
        """完成回复"""
        if self.current_ai_widget:
            self.current_ai_widget.configure(state="normal")
            self.current_ai_widget.delete("1.0", END)
            self.current_ai_widget.insert("1.0", full_response)
            self.current_ai_widget.configure(state=DISABLED)

            time_str = datetime.now().strftime("%H:%M")
            ttkb.Label(
                self.current_ai_container, text=time_str,
                font=("Microsoft YaHei UI", 8),
                bootstyle="secondary"
            ).pack(anchor="w", padx=5, pady=(2, 0))

        self.messages.append({"role": "assistant", "content": full_response})
        self._reset_state()

    def _handle_error(self, msg):
        """处理错误"""
        if self.current_ai_widget:
            self.current_ai_widget.configure(state="normal")
            self.current_ai_widget.delete("1.0", END)
            self.current_ai_widget.insert("1.0", f"❌ {msg}")
            self.current_ai_widget.configure(state=DISABLED)

            time_str = datetime.now().strftime("%H:%M")
            ttkb.Label(
                self.current_ai_container, text=time_str,
                font=("Microsoft YaHei UI", 8),
                bootstyle="secondary"
            ).pack(anchor="w", padx=5, pady=(2, 0))

        self._reset_state()

    def _reset_state(self):
        """重置状态"""
        self.is_loading = False
        self.send_btn.configure(text="发送 ▶", state=NORMAL)
        self.current_ai_widget = None
        self.current_ai_container = None
        self._update_status("就绪", "success")

    def _update_status(self, text, style="secondary"):
        self.status_label.configure(text=f"● {text}", bootstyle=style)

    def _scroll_to_bottom(self):
        self.chat_canvas.yview_moveto(1.0)
        self.root.update_idletasks()

    def _on_enter_key(self, event):
        if event.state & 0x1:  # Shift 按下
            return  # 插入换行
        self.send_message()
        return "break"

    # ════════════════════════════════════════════════════════════════
    #  设置窗口
    # ════════════════════════════════════════════════════════════════

    def open_settings(self):
        win = ttkb.Toplevel(self.root)
        win.title("设置 - DeepSeek 学习助手")
        win.geometry("480x360")
        win.resizable(False, False)
        win.transient(self.root)
        win.grab_set()

        main = ttkb.Frame(win, padding=20)
        main.pack(fill=BOTH, expand=True)

        ttkb.Label(main, text="⚙ 设置",
            font=("Microsoft YaHei UI", 14, "bold")).pack(anchor="w", pady=(0, 15))

        # API Key
        ttkb.Label(main, text="DeepSeek API Key",
            font=("Microsoft YaHei UI", 10)).pack(anchor="w")
        entry = ttkb.Entry(main, font=("Consolas", 10), show="*")
        entry.pack(fill=X, pady=(4, 3))
        entry.insert(0, self.config.get("api_key", ""))

        ttkb.Label(main, text="👉 platform.deepseek.com 获取",
            font=("Microsoft YaHei UI", 8),
            bootstyle="info").pack(anchor="w", pady=(0, 10))

        # 模型
        ttkb.Label(main, text="模型选择",
            font=("Microsoft YaHei UI", 10)).pack(anchor="w")
        combo = ttkb.Combobox(
            main, values=["deepseek-v4-flash", "deepseek-v4-pro", "deepseek-chat (旧)", "deepseek-reasoner (旧)"],
            state="readonly", font=("Microsoft YaHei UI", 10))
        combo.pack(fill=X, pady=(4, 3))
        combo.set(self.config.get("model", DEFAULT_MODEL))

        ttkb.Label(main, text="deepseek-v4-flash: 最新快速版 ⭐ | deepseek-v4-pro: 旗舰版\ndeepseek-chat/reasoner (旧): 2026-07-24 停用",
            font=("Microsoft YaHei UI", 8),
            bootstyle="secondary").pack(anchor="w", pady=(0, 15))

        # 按钮
        btn_f = ttkb.Frame(main)
        btn_f.pack(fill=X)
        ttkb.Button(btn_f, text="保存", bootstyle="success",
            command=lambda: self._save_settings(win, entry.get().strip(), combo.get()),
            width=12).pack(side=RIGHT, padx=(5, 0))
        ttkb.Button(btn_f, text="取消", bootstyle="secondary-outline",
            command=win.destroy, width=12).pack(side=RIGHT)

        win.update_idletasks()
        x = self.root.winfo_x() + (self.root.winfo_width() - win.winfo_width()) // 2
        y = self.root.winfo_y() + (self.root.winfo_height() - win.winfo_height()) // 2
        win.geometry(f"+{x}+{y}")

    def _save_settings(self, win, api_key, model):
        if not api_key:
            messagebox.showwarning("提示", "请输入 DeepSeek API Key")
            return
        self.config["api_key"] = api_key
        self.config["model"] = model
        save_config(self.config)
        self.client = DeepSeekClient(api_key, model)
        self._update_status("已连接", "success")
        win.destroy()

    def _on_close(self):
        self.root.destroy()

    def run(self):
        self.root.mainloop()


# ══════════════════════════════════════════════════════════════════════
#  入口
# ══════════════════════════════════════════════════════════════════════

if __name__ == "__main__":
    # 尝试安装 PyMuPDF（PDF 支持）
    try:
        import fitz
    except ImportError:
        try:
            import pdfminer
        except ImportError:
            pass  # 没有 PDF 支持也可以用
    app = DeepSeekLearnerApp()
    app.run()
