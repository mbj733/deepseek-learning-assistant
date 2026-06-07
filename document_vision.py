#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Document Vision Processor — 文档图像理解模块（国内可用版）

解决 DeepSeek 非多模态的问题，为上传的学习资料提供图像理解能力。

## 国内可用的方案（按推荐顺序）

### 方案 A：通义千问 VL API ⭐（推荐）
  - 阿里云 DashScope，国内直接访问
  - 能理解图表、公式、流程图，不只是 OCR
  - 价格：约 1.6元/千tokens（qwen-vl-plus），极低
  - 获取 Key：https://dashscope.aliyun.com/
  - 设置：在应用设置中填入 DASHSCOPE_API_KEY

### 方案 B：PaddleOCR（本地方案）
  - 百度出品，纯本地免费离线
  - 仅识别图片中的文字（不理解图表含义）
  - 安装：pip install paddleocr
  - 识别中文效果极好

### 方案 C：Ollama 本地视觉模型（本地方案）
  - 完全免费、离线、隐私保护
  - 需要安装 Ollama + 下载视觉模型（如 qwen-vl）
  - CPU 可运行但较慢
"""

import base64
import io
import json
import logging
import os
import re
import tempfile
from pathlib import Path
from typing import Optional, Callable

# 基础图像处理
try:
    from PIL import Image
    HAS_PIL = True
except ImportError:
    HAS_PIL = False

# PDF 截图
try:
    import fitz  # PyMuPDF
    HAS_FITZ = True
except ImportError:
    HAS_FITZ = False

logger = logging.getLogger(__name__)

# ══════════════════════════════════════════════════════════════════════
#  图像提取器 — 从文档中提取图片
# ══════════════════════════════════════════════════════════════════════

class ImageExtractor:
    """从 PDF/Word/PPT 中提取图片"""

    @staticmethod
    def extract_from_pdf(filepath: str, max_pages: int = 20,
                         dpi: int = 150,
                         progress_cb: Optional[Callable] = None) -> list[dict]:
        """将 PDF 每页渲染为图片，返回 [{page, image_bytes, desc}, ...]"""
        if not HAS_FITZ:
            return []

        results = []
        try:
            doc = fitz.open(filepath)
            total = min(len(doc), max_pages)

            for i in range(total):
                page = doc[i]
                pix = page.get_pixmap(dpi=dpi)
                img_bytes = pix.tobytes("png")
                results.append({
                    "page": i + 1,
                    "source": f"第{i+1}页",
                    "image_bytes": img_bytes,
                    "description": "",
                })
                if progress_cb:
                    progress_cb(i + 1, total)

            doc.close()
        except Exception as e:
            logger.warning(f"PDF 截图失败: {e}")

        return results

    @staticmethod
    def extract_from_docx(filepath: str) -> list[dict]:
        """从 Word 文档中提取内嵌图片"""
        results = []
        try:
            from docx import Document
            doc = Document(filepath)
            img_counter = 0
            for rel in doc.part.rels.values():
                if "image" in rel.reltype:
                    img_counter += 1
                    image_data = rel.target_part.blob
                    results.append({
                        "page": img_counter,
                        "source": f"图片{img_counter}",
                        "image_bytes": image_data,
                        "description": "",
                    })
            return results
        except Exception as e:
            logger.warning(f"Word 图片提取失败: {e}")
            return []

    @staticmethod
    def extract_from_pptx(filepath: str) -> list[dict]:
        """从 PPT 中提取内嵌图片"""
        results = []
        try:
            from pptx import Presentation
            prs = Presentation(filepath)
            img_counter = 0

            for slide_num, slide in enumerate(prs.slides, 1):
                for shape in slide.shapes:
                    if shape.shape_type == 13:  # Picture
                        img_counter += 1
                        image = shape.image
                        img_bytes = image.blob
                        results.append({
                            "page": slide_num,
                            "source": f"第{slide_num}页图片",
                            "image_bytes": img_bytes,
                            "description": "",
                        })
            return results
        except Exception as e:
            logger.warning(f"PPT 图片提取失败: {e}")
            return []

    @staticmethod
    def extract(filepath: str, **kwargs) -> list[dict]:
        """自动识别文档类型并提取图片"""
        ext = os.path.splitext(filepath)[1].lower()

        if ext == ".pdf":
            return ImageExtractor.extract_from_pdf(filepath, **kwargs)
        elif ext == ".docx":
            return ImageExtractor.extract_from_docx(filepath)
        elif ext == ".pptx":
            return ImageExtractor.extract_from_pptx(filepath)
        return []


# ══════════════════════════════════════════════════════════════════════
#  图像理解后端
# ══════════════════════════════════════════════════════════════════════

class DashScopeVision:
    """阿里云通义千问 VL API（国内可用，推荐 ⭐）

    优势：
      - 国内直接访问，无需翻墙
      - 理解图表、公式、截图中的文字
      - 价格极低（qwen-vl-plus 约 1.6元/千tokens）
      - 支持中文描述

    获取 Key: https://dashscope.aliyun.com/

    用法：
      engine = DashScopeVision(api_key="sk-xxx")
      desc = engine.describe(image_bytes)
    """

    def __init__(self, api_key: str = ""):
        self.api_key = api_key
        self.model = "qwen-vl-plus"  # 性价比最高

    def is_available(self) -> bool:
        return bool(self.api_key)

    def describe(self, image_bytes: bytes, prompt: str = "") -> str:
        """描述一张图片"""
        if not self.api_key:
            return ""

        try:
            import requests

            # base64 编码图片
            img_b64 = base64.b64encode(image_bytes).decode("utf-8")

            user_prompt = prompt or (
                "请详细描述这张图片中的内容。"
                "如果是文字/表格，请完整提取所有文字；"
                "如果是图表/流程图，请解释其结构和含义；"
                "如果是公式，请用 LaTeX 格式表示。"
                "回答用中文。"
            )

            # DashScope 兼容 OpenAI 格式
            url = "https://dashscope.aliyuncs.com/compatible-mode/v1/chat/completions"
            headers = {
                "Authorization": f"Bearer {self.api_key}",
                "Content-Type": "application/json"
            }
            payload = {
                "model": self.model,
                "messages": [
                    {
                        "role": "user",
                        "content": [
                            {"type": "image_url", "image_url": f"data:image/png;base64,{img_b64}"},
                            {"type": "text", "text": user_prompt}
                        ]
                    }
                ],
                "max_tokens": 512,
                "temperature": 0.1,
            }

            resp = requests.post(url, headers=headers, json=payload, timeout=60)
            resp.raise_for_status()
            data = resp.json()

            text = data.get("choices", [{}])[0].get("message", {}).get("content", "")
            return text.strip() if text else "[无法描述此图片]"

        except Exception as e:
            return f"[图片描述失败: {e}]"


class PaddleOCREngine:
    """PaddleOCR 本地 OCR（国内可用，免费离线）

    安装：pip install paddleocr

    优势：
      - 纯本地，免费，离线
      - 中文识别效果极好
      - 但只能识别文字，不理解图表含义

    兼容 PaddleOCR 2.x / 3.x
    """

    def __init__(self):
        self._ocr = None
        self._available = False
        self._is_v3 = False  # PaddleOCR 3.x 标记
        self._init()

    def _init(self):
        try:
            # PaddleOCR 3.x 导入路径
            import os
            os.environ['PADDLE_PDX_DISABLE_MODEL_SOURCE_CHECK'] = 'True'

            try:
                from paddleocr._pipelines import PaddleOCR as _PaddleOCR
                self._ocr = _PaddleOCR()
                self._is_v3 = True
            except (ImportError, AttributeError):
                try:
                    from paddleocr import PaddleOCR as _PaddleOCR
                    self._ocr = _PaddleOCR()
                    self._is_v3 = True
                except (ImportError, TypeError, ValueError):
                    # PaddleOCR 2.x 旧版 API
                    from paddleocr import PaddleOCR as _PaddleOCR
                    self._ocr = _PaddleOCR(use_angle_cls=True, lang='ch',
                                          show_log=False, use_gpu=False)
                    self._is_v3 = False

            self._available = True

        except ImportError:
            self._available = False
        except Exception as e:
            logger.warning(f"PaddleOCR 初始化失败: {e}")
            self._available = False

    def is_available(self) -> bool:
        return self._available

    def describe(self, image_bytes: bytes, prompt: str = "") -> str:
        """OCR 识别图片中的文字"""
        if not self._available or not self._ocr:
            return ""

        try:
            # 写入临时文件（先 close 再使用，避免 Windows 文件锁）
            import uuid
            tmp_path = os.path.join(tempfile.gettempdir(),
                                    f"ocr_{uuid.uuid4().hex}.png")
            with open(tmp_path, "wb") as f:
                f.write(image_bytes)

            # 调用 OCR（根据版本选择 API）
            if self._is_v3:
                result = self._ocr.ocr(tmp_path)
            else:
                result = self._ocr.ocr(tmp_path, cls=True)

            # 清理临时文件
            try:
                os.unlink(tmp_path)
            except:
                pass

            # 解析结果
            if not result:
                return "[图片中未识别到文字]"

            # PaddleOCR 3.x 返回格式可能不同
            texts = []
            if isinstance(result, list) and len(result) > 0:
                lines = result[0] if isinstance(result[0], list) else result
                for line in lines:
                    if isinstance(line, (list, tuple)) and len(line) >= 2:
                        text_info = line[1] if isinstance(line[1], (list, tuple)) else line
                        if isinstance(text_info, (list, tuple)) and len(text_info) >= 2:
                            text = str(text_info[0])
                            confidence = float(text_info[1])
                            if confidence > 0.3:
                                texts.append(text)

            if texts:
                return "【图片中的文字】\n" + "\n".join(texts)
            return "[图片中未识别到清晰文字]"

        except Exception as e:
            return f"[OCR 失败: {e}]"


class OllamaVision:
    """Ollama 本地视觉模型（国内可用，完全免费离线）

    要求：
      1. 安装 Ollama: https://ollama.com/
      2. 下载视觉模型: ollama pull qwen2.5-vl:7b
         (或 minicpm-v, llava 等)

    优势：
      - 完全免费，完全离线，隐私保护
      - CPU 可运行（较慢），GPU 更快

    用法：
      engine = OllamaVision(model="qwen2.5-vl:7b")
      desc = engine.describe(image_bytes)
    """

    def __init__(self, model: str = "qwen2.5-vl:7b", base_url: str = "http://localhost:11434"):
        self.model = model
        self.base_url = base_url.rstrip("/")
        self._available = self._check()

    def _check(self) -> bool:
        try:
            import requests
            resp = requests.get(f"{self.base_url}/api/tags", timeout=5)
            if resp.status_code == 200:
                models = [m["name"] for m in resp.json().get("models", [])]
                # 检查是否有视觉模型
                for m in models:
                    if any(v in m.lower() for v in ["vl", "vision", "llava", "minicpm"]):
                        self.model = m  # 使用找到的第一个视觉模型
                        return True
                # 没有找到视觉模型，但 Ollama 在运行
                return False
            return False
        except:
            return False

    def is_available(self) -> bool:
        return self._available

    def describe(self, image_bytes: bytes, prompt: str = "") -> str:
        if not self._available:
            return ""

        try:
            import requests

            img_b64 = base64.b64encode(image_bytes).decode("utf-8")
            user_prompt = prompt or "请详细描述这张图片中的内容，包括所有文字、图表、公式。"

            payload = {
                "model": self.model,
                "messages": [
                    {
                        "role": "user",
                        "content": user_prompt,
                        "images": [img_b64]
                    }
                ],
                "stream": False,
                "options": {"temperature": 0.1}
            }

            resp = requests.post(
                f"{self.base_url}/api/chat",
                json=payload, timeout=120
            )
            resp.raise_for_status()
            data = resp.json()
            content = data.get("message", {}).get("content", "")
            return content.strip() if content else "[无法描述此图片]"

        except Exception as e:
            return f"[本地模型描述失败: {e}]"


class GeminiVision:
    """Google Gemini API（国内不可用，保留兼容）

    注意：国内网络无法访问 Google API，
    请使用 DashScopeVision 或 PaddleOCREngine 代替。
    """

    BASE_URL = "https://generativelanguage.googleapis.com/v1beta/models"

    def __init__(self, api_key: str = ""):
        self.api_key = api_key
        self.model = "gemini-2.0-flash-lite"

    def is_available(self) -> bool:
        return bool(self.api_key)

    def describe(self, image_bytes: bytes, prompt: str = "") -> str:
        if not self.api_key:
            return ""
        try:
            import requests
            img_b64 = base64.b64encode(image_bytes).decode("utf-8")
            mime = "image/png"
            url = f"{self.BASE_URL}/{self.model}:generateContent?key={self.api_key}"
            user_prompt = prompt or "请详细描述这张图片中的内容，包括所有文字、图表、公式。"

            payload = {
                "contents": [{
                    "parts": [
                        {"inline_data": {"mime_type": mime, "data": img_b64}},
                        {"text": user_prompt}
                    ]
                }],
                "generationConfig": {"temperature": 0.1, "maxOutputTokens": 512}
            }

            resp = requests.post(url, json=payload, timeout=30)
            resp.raise_for_status()
            data = resp.json()
            candidates = data.get("candidates", [])
            if candidates:
                text = candidates[0].get("content", {}).get("parts", [{}])[0].get("text", "")
                return text.strip()
            return "[Gemini: 无法描述此图片]"
        except Exception as e:
            return f"[图片描述失败: {e}]"


# ══════════════════════════════════════════════════════════════════════
#  统一处理器 — 主入口
# ══════════════════════════════════════════════════════════════════════

class DocumentVisionProcessor:
    """文档图像处理器 — 主入口

    支持的 backend:
      - "dashscope":  阿里通义千问 VL API（国内推荐 ⭐）
      - "paddleocr":  PaddleOCR 本地方案
      - "ollama":     Ollama 本地视觉模型
      - "gemini":     Google Gemini（国内不可用）

    用法：
      # 方案 A: 通义千问（国内推荐）
      vp = DocumentVisionProcessor(
          backend="dashscope",
          api_key="sk-xxxx"  # 从 dashscope.aliyun.com 获取
      )
      desc = vp.process_document("教材.pdf")

      # 方案 B: PaddleOCR 本地
      # pip install paddleocr
      vp = DocumentVisionProcessor(backend="paddleocr")
      desc = vp.process_document("教材.pdf")
    """

    BACKENDS = {
        "dashscope": DashScopeVision,
        "paddleocr": PaddleOCREngine,
        "ollama": OllamaVision,
        "gemini": GeminiVision,
    }

    BACKEND_NAMES = {
        "dashscope": "通义千问 VL API（推荐 ⭐）",
        "paddleocr": "PaddleOCR 本地文字识别",
        "ollama": "Ollama 本地视觉模型",
        "gemini": "Gemini API（国内不可用）",
    }

    def __init__(self, backend: str = "dashscope", api_key: str = "",
                 ollama_model: str = "qwen2.5-vl:7b",
                 describe_prompt: str = ""):
        """
        Args:
            backend: "dashscope" | "paddleocr" | "ollama" | "gemini"
            api_key: DashScope 或 Gemini 的 API Key
            ollama_model: Ollama 使用的模型名
            describe_prompt: 自定义图片描述提示词
        """
        self.backend_name = backend
        self.describe_prompt = describe_prompt or (
            "请详细描述这张图片中的内容。"
            "如果是文字/表格，请完整提取所有文字内容；"
            "如果是图表/流程图，请解释其结构和含义；"
            "如果是公式，请用 LaTeX 格式表示。"
            "回答请用中文。"
        )
        self.ollama_model = ollama_model
        self._engine = None
        self._init_engine(api_key)

    def _init_engine(self, api_key: str):
        if self.backend_name == "dashscope":
            self._engine = DashScopeVision(api_key)
        elif self.backend_name == "paddleocr":
            self._engine = PaddleOCREngine()
        elif self.backend_name == "ollama":
            self._engine = OllamaVision(model=self.ollama_model)
        elif self.backend_name == "gemini":
            self._engine = GeminiVision(api_key)
        else:
            self._engine = None

    def is_available(self) -> bool:
        return self._engine is not None and self._engine.is_available()

    def get_backend_name(self) -> str:
        """返回当前后端名称"""
        engine_type = type(self._engine).__name__ if self._engine else "None"
        return self.BACKEND_NAMES.get(self.backend_name, engine_type)

    def process_document(self, filepath: str,
                         progress_cb: Optional[Callable] = None) -> str:
        """处理一个文档文件，返回所有图片描述的合并文本"""
        if not self.is_available():
            return ""

        images = ImageExtractor.extract(filepath, progress_cb=progress_cb)
        if not images:
            return ""

        total = len(images)
        descriptions = []

        for idx, img in enumerate(images):
            desc = self._engine.describe(img["image_bytes"], self.describe_prompt)
            if desc and not desc.startswith("[图片"):
                descriptions.append(f"【{img['source']}】\n{desc}")

            if progress_cb:
                progress_cb(idx + 1, total)

        if descriptions:
            return "\n\n" + "\n\n".join(descriptions)
        return ""


# ══════════════════════════════════════════════════════════════════════
#  快速测试
# ══════════════════════════════════════════════════════════════════════

if __name__ == "__main__":
    print("=" * 50)
    print("📷 Document Vision Processor 测试")
    print("=" * 50)

    print(f"\n✅ PyMuPDF (PDF截图): {'可用' if HAS_FITZ else '不可用'}")
    print(f"✅ PIL (图像处理):     {'可用' if HAS_PIL else '不可用'}")

    # 检测可用后端
    print("\n📡 检测可用后端:")

    dash_key = os.environ.get("DASHSCOPE_API_KEY", "")
    if dash_key:
        d = DashScopeVision(dash_key)
        print(f"  ✅ 通义千问 VL API: 已配置")
    else:
        print(f"  💡 通义千问 VL API: 未配置（推荐！）")

    ocr = PaddleOCREngine()
    print(f"  {'✅' if ocr.is_available() else '❌'} PaddleOCR: {'可用' if ocr.is_available() else '未安装 (pip install paddleocr)'}")

    ollama = OllamaVision()
    print(f"  {'✅' if ollama.is_available() else '💡'} Ollama 本地模型: {'可用' if ollama.is_available() else '未安装 (ollama pull qwen2.5-vl:7b)'}")

    print(f"\n{'='*50}")
    print(f"💡 国内推荐方案:")
    print(f"   1️⃣  通义千问 VL API（在线，理解图表+文字）")
    print(f"       Key: https://dashscope.aliyun.com/")
    print(f"   2️⃣  PaddleOCR（本地免费，仅文字识别）")
    print(f"       安装: pip install paddleocr")
    print(f"   3️⃣  Ollama（本地免费，需下载模型）")
    print(f"       安装: ollama pull qwen2.5-vl:7b")
