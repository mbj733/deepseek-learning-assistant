#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
DeepSeek 学习助手 v3.0
基于 Hermes Agent 架构模式重构

核心特性：
  - SQLite + FTS5 全文搜索知识库（借鉴 hermes_state.py）
  - 学习卡片系统（借鉴 skills 系统）
  - 三层系统提示词（stable + context + volatile）
  - 会话持久化存储
  - 智能分块与检索排序
"""

import json
import logging
import os
import re
import sqlite3
import sys
import threading
import tkinter as tk
from tkinter import ttk, messagebox, filedialog, scrolledtext, simpledialog
from datetime import datetime
from pathlib import Path
import shutil
from dataclasses import dataclass, field
from typing import List, Optional
import hashlib

import requests
import ttkbootstrap as ttkb
from ttkbootstrap.constants import *

# 文档视觉处理（可选，无依赖时静默降级）
try:
    from document_vision import DocumentVisionProcessor, ImageExtractor
    HAS_VISION = True
except ImportError:
    HAS_VISION = False
    DocumentVisionProcessor = None

# ══════════════════════════════════════════════════════════════════════
#  路径与常量 — 支持 PyInstaller 打包（exe 也能持久化保存）
# ══════════════════════════════════════════════════════════════════════

if getattr(sys, 'frozen', False):
    # 打包成 .exe 运行时，配置文件保存在 exe 同目录
    APP_DIR = os.path.dirname(os.path.abspath(sys.executable))
else:
    # 源码运行时，保存在脚本所在目录
    APP_DIR = os.path.dirname(os.path.abspath(__file__))

DB_PATH = os.path.join(APP_DIR, "sessions.db")
SESSIONS_DIR = os.path.join(APP_DIR, "sessions")
CONFIG_FILE = os.path.join(APP_DIR, "config.yaml")

DEEPSEEK_API_URL = "https://api.deepseek.com/chat/completions"
DEFAULT_MODEL = "deepseek-v4-flash"

CHUNK_SIZE = 600       # 知识块大小
CHUNK_OVERLAP = 150    # 重叠字符

os.makedirs(SESSIONS_DIR, exist_ok=True)

# ══════════════════════════════════════════════════════════════════════
#  SQLite 数据库 — 借鉴 hermes_state.py 的 WAL + FTS5 设计
# ══════════════════════════════════════════════════════════════════════

SCHEMA_SQL = """
PRAGMA journal_mode=WAL;
PRAGMA foreign_keys=ON;

CREATE TABLE IF NOT EXISTS sessions (
    id TEXT PRIMARY KEY,
    name TEXT NOT NULL DEFAULT '未命名',
    created_at TEXT NOT NULL,
    updated_at TEXT NOT NULL,
    message_count INTEGER DEFAULT 0,
    card_count INTEGER DEFAULT 0
);

CREATE TABLE IF NOT EXISTS messages (
    id INTEGER PRIMARY KEY AUTOINCREMENT,
    session_id TEXT NOT NULL REFERENCES sessions(id) ON DELETE CASCADE,
    role TEXT NOT NULL,
    content TEXT NOT NULL,
    created_at TEXT NOT NULL,
    token_count INTEGER DEFAULT 0
);

CREATE VIRTUAL TABLE IF NOT EXISTS messages_fts USING fts5(
    content, tokenize='unicode61'
);

CREATE TABLE IF NOT EXISTS knowledge_chunks (
    id INTEGER PRIMARY KEY AUTOINCREMENT,
    session_id TEXT NOT NULL REFERENCES sessions(id) ON DELETE CASCADE,
    source_file TEXT NOT NULL,
    chunk_index INTEGER NOT NULL,
    content TEXT NOT NULL,
    embedding BLOB
);

CREATE VIRTUAL TABLE IF NOT EXISTS knowledge_chunks_fts USING fts5(
    content, tokenize='unicode61'
);

CREATE TABLE IF NOT EXISTS knowledge_files (
    id INTEGER PRIMARY KEY AUTOINCREMENT,
    session_id TEXT NOT NULL REFERENCES sessions(id) ON DELETE CASCADE,
    file_name TEXT NOT NULL,
    file_path TEXT NOT NULL,
    file_size INTEGER DEFAULT 0,
    chunk_count INTEGER DEFAULT 0,
    added_at TEXT NOT NULL
);

-- 学习卡片系统（借鉴 Hermes skills）
CREATE TABLE IF NOT EXISTS study_cards (
    id INTEGER PRIMARY KEY AUTOINCREMENT,
    session_id TEXT NOT NULL REFERENCES sessions(id) ON DELETE CASCADE,
    title TEXT NOT NULL,
    question TEXT,
    answer TEXT NOT NULL,
    category TEXT DEFAULT '通用',
    tags TEXT DEFAULT '',
    created_at TEXT NOT NULL,
    review_count INTEGER DEFAULT 0,
    difficulty INTEGER DEFAULT 3
);

-- 触发器：同步 FTS
CREATE TRIGGER IF NOT EXISTS knowledge_chunks_ai AFTER INSERT ON knowledge_chunks BEGIN
    INSERT INTO knowledge_chunks_fts(rowid, content) VALUES (new.id, new.content);
END;

CREATE TRIGGER IF NOT EXISTS knowledge_chunks_ad AFTER DELETE ON knowledge_chunks BEGIN
    INSERT INTO knowledge_chunks_fts(knowledge_chunks_fts, rowid, content) VALUES('delete', old.id, old.content);
END;

CREATE TRIGGER IF NOT EXISTS knowledge_chunks_au AFTER UPDATE ON knowledge_chunks BEGIN
    INSERT INTO knowledge_chunks_fts(knowledge_chunks_fts, rowid, content) VALUES('delete', old.id, old.content);
    INSERT INTO knowledge_chunks_fts(rowid, content) VALUES (new.id, new.content);
END;
"""


class Database:
    """SQLite 数据库管理（借鉴 Hermes state.db 设计）"""

    _instance = None
    _lock = threading.Lock()

    def __new__(cls):
        if cls._instance is None:
            with cls._lock:
                if cls._instance is None:
                    cls._instance = super().__new__(cls)
                    cls._instance._init_db()
        return cls._instance

    def _init_db(self):
        self.conn = sqlite3.connect(DB_PATH, check_same_thread=False)
        self.conn.row_factory = sqlite3.Row
        self.conn.executescript(SCHEMA_SQL)
        self.conn.commit()
        self._local = threading.local()

    def _get_conn(self):
        """线程安全的连接"""
        if not hasattr(self._local, 'conn') or self._local.conn is None:
            self._local.conn = sqlite3.connect(DB_PATH, check_same_thread=False)
            self._local.conn.row_factory = sqlite3.Row
        return self._local.conn

    def execute(self, sql, params=None):
        return self.conn.execute(sql, params or [])

    def fetchall(self, sql, params=None):
        return self.conn.execute(sql, params or []).fetchall()

    def fetchone(self, sql, params=None):
        return self.conn.execute(sql, params or []).fetchone()

    def commit(self):
        self.conn.commit()


# ══════════════════════════════════════════════════════════════════════
#  知识库引擎 — FTS5 支持（借鉴 Hermes RAG 模式）
# ══════════════════════════════════════════════════════════════════════

class KnowledgeEngine:
    """知识库引擎 — 基于 SQLite FTS5 的检索增强生成"""

    def __init__(self, session_id: str, vision_processor=None):
        self.session_id = session_id
        self.db = Database()
        self.vision = vision_processor  # DocumentVisionProcessor 实例

    def add_file(self, src_path: str) -> tuple[bool, str]:
        """添加文件到知识库，自动分块、OCR图片并建立 FTS 索引"""
        name = os.path.basename(src_path)
        base, ext = os.path.splitext(name)
        files_dir = os.path.join(SESSIONS_DIR, self.session_id, "files")
        os.makedirs(files_dir, exist_ok=True)

        # 避免重名
        dest_path = os.path.join(files_dir, name)
        counter = 1
        while os.path.exists(dest_path):
            name = f"{base}_{counter}{ext}"
            dest_path = os.path.join(files_dir, name)
            counter += 1

        shutil.copy2(src_path, dest_path)

        # 提取文本
        content = self._extract_text(dest_path)

        # 提取并处理图片（如果视觉处理器可用）
        vision_text = ""
        if self.vision and self.vision.is_available():
            try:
                vision_text = self.vision.process_document(dest_path)
            except Exception as e:
                logger.warning(f"图片处理失败: {e}")

        # 将图片描述附加到文本末尾
        if vision_text:
            content += "\n\n" + ("=" * 40) + "\n【文档图片内容】\n" + vision_text
        if not content or len(content.strip()) < 20:
            os.remove(dest_path)
            return False, f"无法读取「{name}」的内容"

        # 分块
        chunks = self._chunk_text(content)

        # 存入数据库
        file_size = os.path.getsize(dest_path)
        now = datetime.now().isoformat()

        self.db.execute(
            """INSERT INTO knowledge_files (session_id, file_name, file_path, file_size, chunk_count, added_at)
               VALUES (?, ?, ?, ?, ?, ?)""",
            (self.session_id, name, dest_path, file_size, len(chunks), now)
        )

        for i, chunk_text in enumerate(chunks):
            self.db.execute(
                """INSERT INTO knowledge_chunks (session_id, source_file, chunk_index, content)
                   VALUES (?, ?, ?, ?)""",
                (self.session_id, name, i, chunk_text)
            )

        self.db.commit()
        return True, f"✅ 已导入「{name}」（{len(chunks)} 个片段）"

    def remove_file(self, file_name: str) -> tuple[bool, str]:
        """删除文件及其所有分块"""
        rows = self.db.fetchall(
            "SELECT id, file_path FROM knowledge_files WHERE session_id=? AND file_name=?",
            (self.session_id, file_name)
        )
        if not rows:
            return False, "文件不存在"

        # 删除物理文件
        for row in rows:
            if os.path.exists(row["file_path"]):
                os.remove(row["file_path"])

        # 删除数据库记录（级联删除 chunks）
        self.db.execute(
            "DELETE FROM knowledge_chunks WHERE session_id=? AND source_file=?",
            (self.session_id, file_name)
        )
        self.db.execute(
            "DELETE FROM knowledge_files WHERE session_id=? AND file_name=?",
            (self.session_id, file_name)
        )
        self.db.commit()
        return True, f"已移除「{file_name}」"

    def search(self, query: str, top_k: int = 5) -> list[dict]:
        """FTS5 全文搜索 + BM25 排序"""
        if not query.strip():
            return []

        try:
            # FTS5 精确匹配
            fts_query = self._build_fts_query(query)
            rows = self.db.fetchall("""
                SELECT kc.content, kc.source_file, kc.chunk_index,
                       rank_bm25(knowledge_chunks_fts) as score
                FROM knowledge_chunks_fts
                JOIN knowledge_chunks kc ON knowledge_chunks_fts.rowid = kc.id
                WHERE knowledge_chunks_fts MATCH ?
                  AND kc.session_id = ?
                ORDER BY score
                LIMIT ?
            """, (fts_query, self.session_id, top_k * 2))

            results = []
            seen_sources = set()
            for row in rows:
                source_key = row["source_file"] + ":" + str(row["chunk_index"])
                if source_key not in seen_sources:
                    seen_sources.add(source_key)
                    results.append({
                        "text": row["content"],
                        "source": row["source_file"],
                        "chunk": row["chunk_index"],
                        "score": row["score"]
                    })

            return results[:top_k]

        except sqlite3.OperationalError:
            # FTS5 查询语法错误时回退到 LIKE 搜索
            return self._fallback_search(query, top_k)

    def _build_fts_query(self, query: str) -> str:
        """将自然语言查询转换为 FTS5 查询"""
        # 去除标点，保留中文字符和单词
        words = re.findall(r'[\w\u4e00-\u9fff]+', query)
        fts_terms = []
        for w in words:
            if len(w) <= 1:
                continue
            if re.match(r'^[\u4e00-\u9fff]+$', w):
                # 中文：拆成单字加前缀匹配
                fts_terms.append(f'"{w}"')
            else:
                fts_terms.append(w)

        if not fts_terms:
            # 如果都无法匹配，用 OR 连接单个字
            chars = [c for c in query if c.strip()]
            fts_terms = [f'"{c}"' for c in chars[:10]]

        return " AND ".join(fts_terms)

    def _fallback_search(self, query: str, top_k: int) -> list[dict]:
        """FTS 不可用时的 LIKE 回退"""
        # 获取该会话所有 chunks
        rows = self.db.fetchall(
            "SELECT content, source_file, chunk_index FROM knowledge_chunks WHERE session_id=?",
            (self.session_id,)
        )

        query_lower = query.lower()
        query_words = set(re.findall(r'[\w\u4e00-\u9fff]+', query_lower))

        scored = []
        for row in rows:
            text_lower = row["content"].lower()
            text_words = set(re.findall(r'[\w\u4e00-\u9fff]+', text_lower))
            common = query_words & text_words
            if not common:
                continue

            score = len(common)
            for qw in query_words:
                if len(qw) > 1:
                    score += min(text_lower.count(qw), 10) * 0.3

            scored.append((score, row))

        scored.sort(key=lambda x: -x[0])
        top = scored[:top_k]

        results = []
        seen = set()
        for _, row in top:
            key = row["source_file"] + ":" + str(row["chunk_index"])
            if key not in seen:
                seen.add(key)
                results.append({
                    "text": row["content"],
                    "source": row["source_file"],
                    "chunk": row["chunk_index"],
                    "score": _,
                })
        return results

    def format_context(self, query: str) -> str:
        """格式化搜索结果为上下文文本"""
        results = self.search(query)
        if not results:
            return "（当前资料库中未找到相关内容）"

        parts = []
        for i, r in enumerate(results, 1):
            source = r["source"]
            parts.append(f"[参考资料 {i} — 来自《{source}》]\n{r['text']}\n")

        return "\n\n---\n\n".join(parts)

    def get_stats(self) -> tuple:
        """获取知识库统计"""
        files = self.db.fetchall(
            "SELECT COUNT(*) as cnt, SUM(chunk_count) as chunks FROM knowledge_files WHERE session_id=?",
            (self.session_id,)
        )
        row = files[0] if files else {"cnt": 0, "chunks": 0}
        total_chars = self.db.fetchone(
            "SELECT SUM(LENGTH(content)) as chars FROM knowledge_chunks WHERE session_id=?",
            (self.session_id,)
        )
        chars = total_chars["chars"] or 0 if total_chars else 0
        return row["cnt"] or 0, row["chunks"] or 0, chars

    def get_files(self) -> list[dict]:
        """获取文件列表"""
        return [dict(r) for r in self.db.fetchall(
            "SELECT file_name, file_size, chunk_count, added_at FROM knowledge_files WHERE session_id=? ORDER BY added_at DESC",
            (self.session_id,)
        )]

    def _extract_text(self, filepath: str) -> str:
        """提取文本（支持 txt/md/pdf/docx/pptx）"""
        ext = os.path.splitext(filepath)[1].lower()

        if ext in (".txt", ".md"):
            for enc in ("utf-8", "gbk", "utf-16"):
                try:
                    with open(filepath, "r", encoding=enc) as f:
                        return f.read()
                except:
                    continue
            return ""

        elif ext == ".pdf":
            return self._extract_pdf(filepath)
        elif ext == ".docx":
            return self._extract_docx(filepath)
        elif ext == ".pptx":
            return self._extract_pptx(filepath)

        return ""

    def _extract_pdf(self, filepath: str) -> str:
        try:
            import fitz
            doc = fitz.open(filepath)
            text = "\n".join(page.get_text() for page in doc)
            doc.close()
            return text
        except:
            return ""

    def _extract_docx(self, filepath: str) -> str:
        try:
            from docx import Document
            doc = Document(filepath)
            paras = [p.text for p in doc.paragraphs if p.text.strip()]
            tables = []
            for table in doc.tables:
                for row in table.rows:
                    cells = [c.text.strip() for c in row.cells if c.text.strip()]
                    if cells:
                        tables.append(" | ".join(cells))
            result = "\n".join(paras)
            if tables:
                result += "\n\n【表格内容】\n" + "\n".join(tables)
            return result
        except:
            return ""

    def _extract_pptx(self, filepath: str) -> str:
        try:
            from pptx import Presentation
            prs = Presentation(filepath)
            pages = []
            for num, slide in enumerate(prs.slides, 1):
                texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            if para.text.strip():
                                texts.append(para.text.strip())
                    if shape.has_table:
                        for row in shape.table.rows:
                            cells = [c.text.strip() for c in row.cells if c.text.strip()]
                            if cells:
                                texts.append(" | ".join(cells))
                if texts:
                    pages.append(f"【第{num}页】\n" + "\n".join(texts))
            return "\n\n".join(pages)
        except:
            return ""

    def _chunk_text(self, text: str) -> list[str]:
        """智能分块：按段落分割 + 自适应大小"""
        if not text:
            return []

        # 按段落分割
        paragraphs = re.split(r'\n\s*\n', text)
        chunks = []
        current = ""

        for para in paragraphs:
            para = para.strip()
            if not para:
                continue

            if len(current) + len(para) < CHUNK_SIZE:
                current = (current + "\n\n" + para) if current else para
            else:
                if current:
                    chunks.append(current)
                current = self._split_long_para(para) if len(para) > CHUNK_SIZE else para

        if current:
            chunks.append(current)

        return chunks

    def _split_long_para(self, text: str) -> str:
        """按句子分割超长段落"""
        # 尝试按句号、问号、感叹号分割
        sentences = re.split(r'(?<=[。！？.!?])\s*', text)
        if len(sentences) <= 1:
            # 无法分割，直接返回
            return text
        # 只返回前几段保持合理长度
        result = ""
        for s in sentences:
            if len(result) + len(s) < CHUNK_SIZE:
                result += s
        return result.strip() or text[:CHUNK_SIZE]


# ══════════════════════════════════════════════════════════════════════
#  学习卡片系统（借鉴 Hermes Skills）
# ══════════════════════════════════════════════════════════════════════

class StudyCardSystem:
    """学习卡片系统 — 类似于 Hermes 的 Skill 系统"""

    def __init__(self, session_id: str):
        self.session_id = session_id
        self.db = Database()

    def add_card(self, title: str, answer: str, question: str = "",
                 category: str = "通用", tags: str = "") -> int:
        now = datetime.now().isoformat()
        self.db.execute(
            """INSERT INTO study_cards (session_id, title, question, answer, category, tags, created_at)
               VALUES (?, ?, ?, ?, ?, ?, ?)""",
            (self.session_id, title, question, answer, category, tags, now)
        )
        self.db.execute(
            "UPDATE sessions SET card_count = card_count + 1 WHERE id = ?",
            (self.session_id,)
        )
        self.db.commit()
        return self.db.conn.execute("SELECT last_insert_rowid()").fetchone()[0]

    def get_cards(self, category: str = "") -> list[dict]:
        if category:
            rows = self.db.fetchall(
                "SELECT * FROM study_cards WHERE session_id=? AND category=? ORDER BY created_at DESC",
                (self.session_id, category)
            )
        else:
            rows = self.db.fetchall(
                "SELECT * FROM study_cards WHERE session_id=? ORDER BY created_at DESC",
                (self.session_id,)
            )
        return [dict(r) for r in rows]

    def get_categories(self) -> list[str]:
        rows = self.db.fetchall(
            "SELECT DISTINCT category FROM study_cards WHERE session_id=? ORDER BY category",
            (self.session_id,)
        )
        return [r["category"] for r in rows]

    def delete_card(self, card_id: int):
        self.db.execute("DELETE FROM study_cards WHERE id=?", (card_id,))
        self.db.execute(
            "UPDATE sessions SET card_count = MAX(0, card_count - 1) WHERE id = ?",
            (self.session_id,)
        )
        self.db.commit()

    def review_card(self, card_id: int, difficulty: int):
        self.db.execute(
            "UPDATE study_cards SET review_count = review_count + 1, difficulty = ? WHERE id = ?",
            (difficulty, card_id)
        )
        self.db.commit()


# ══════════════════════════════════════════════════════════════════════
#  会话管理（基于 SQLite，借鉴 Hermes SessionDB）
# ══════════════════════════════════════════════════════════════════════

class SessionManager:
    def __init__(self):
        self.db = Database()

    def get_all(self) -> list[dict]:
        rows = self.db.fetchall(
            "SELECT * FROM sessions ORDER BY updated_at DESC"
        )
        return [dict(r) for r in rows]

    def create(self, name: str = "新课程") -> str:
        import uuid
        sid = uuid.uuid4().hex[:12]
        now = datetime.now().isoformat()
        self.db.execute(
            "INSERT INTO sessions (id, name, created_at, updated_at) VALUES (?, ?, ?, ?)",
            (sid, name, now, now)
        )
        self.db.commit()
        # 创建文件目录
        os.makedirs(os.path.join(SESSIONS_DIR, sid, "files"), exist_ok=True)
        return sid

    def rename(self, sid: str, name: str):
        self.db.execute(
            "UPDATE sessions SET name=? WHERE id=?", (name, sid)
        )
        self.db.commit()

    def delete(self, sid: str):
        # 删除物理文件
        sdir = os.path.join(SESSIONS_DIR, sid)
        if os.path.exists(sdir):
            shutil.rmtree(sdir)
        # 删除数据库记录（级联删除自动处理）
        self.db.execute("DELETE FROM sessions WHERE id=?", (sid,))
        self.db.commit()

    def get_name(self, sid: str) -> str:
        row = self.db.fetchone("SELECT name FROM sessions WHERE id=?", (sid,))
        return row["name"] if row else "未命名"

    def update_time(self, sid: str):
        now = datetime.now().isoformat()
        self.db.execute(
            "UPDATE sessions SET updated_at=?, message_count=message_count+1 WHERE id=?",
            (now, sid)
        )
        self.db.commit()


# ══════════════════════════════════════════════════════════════════════
#  AI 客户端
# ══════════════════════════════════════════════════════════════════════

# ── 三层系统提示词（借鉴 Hermes system_prompt.py）──
STABLE_IDENTITY = """你是「DeepSeek 学习助手」，一位专业、耐心的 AI 导师。

【核心原则】
1. 回答必须严格基于用户提供的【参考资料】，不得使用资料外的知识
2. 如果参考资料中找不到答案，明确说"根据现有资料，我无法回答这个问题"
3. 用清晰的语言+生动的类比解释复杂概念
4. 鼓励式教学：先肯定再引导"""

KNOWLEDGE_CONTEXT_TEMPLATE = """
【参考资料】
{context}"""

VOLATILE_NOTE = """
【当前时间】{timestamp}
【学习进度】本次会话已提问 {q_count} 次，已导入 {file_count} 份资料，创建了 {card_count} 张学习卡片
加油！坚持学习一定会进步 💪"""


class DeepSeekClient:
    def __init__(self, api_key: str, model: str = DEFAULT_MODEL):
        self.api_key = api_key
        self.model = model
        self.headers = {
            "Authorization": f"Bearer {api_key}",
            "Content-Type": "application/json"
        }

    def chat_stream(self, messages: list):
        """流式聊天，支持 reasoning_content（deepseek-reasoner 的思维链）"""
        is_reasoning = "reasoner" in self.model

        payload = {
            "model": self.model,
            "messages": messages,
            "stream": True,
            "temperature": 0.3,
            "max_tokens": 4096,
        }
        resp = requests.post(DEEPSEEK_API_URL, headers=self.headers,
                             json=payload, stream=True, timeout=120)
        resp.raise_for_status()

        for line in resp.iter_lines(decode_unicode=True):
            if line:
                if line.startswith("data: "):
                    d = line[6:]
                    if d.strip() == "[DONE]":
                        break
                    try:
                        data = json.loads(d)
                        delta = data["choices"][0].get("delta", {})

                        # Reasoner 模型的思维链
                        reasoning = delta.get("reasoning_content", "")
                        if reasoning:
                            yield ("reasoning", reasoning)
                            continue

                        # 普通内容
                        content = delta.get("content", "")
                        if content:
                            yield ("content", content)

                    except (json.JSONDecodeError, KeyError):
                        continue


# ══════════════════════════════════════════════════════════════════════
#  GUI 应用
# ══════════════════════════════════════════════════════════════════════

FILE_ICONS = {
    ".pdf": "📕", ".txt": "📝", ".md": "📋",
    ".docx": "📘", ".pptx": "📊",
}
CARD_ICONS = {"通用": "📌", "知识点": "🎯", "公式": "🧮",
              "定义": "📖", "代码": "💻", "错题": "❌"}


def _truncate(text: str, max_len: int) -> str:
    """截断长文本，保留扩展名可见"""
    if len(text) <= max_len:
        return text
    # 保留扩展名
    dot = text.rfind(".")
    if dot > 0 and len(text) - dot < 8:
        ext = text[dot:]
        body_max = max_len - len(ext) - 3
        if body_max > 0:
            return text[:body_max] + "…" + ext
    return text[:max_len - 1] + "…"


class DeepSeekLearnerApp:
    def __init__(self):
        self.config = self._load_config()
        self.client = None
        self.sm = SessionManager()
        self.current_sid: Optional[str] = None
        self.kb: Optional[KnowledgeEngine] = None
        self.cards: Optional[StudyCardSystem] = None
        self.vision: Optional = None
        self.conversation: list = []  # 当前会话消息
        self.is_loading = False
        self.current_ai_widget = None
        self.current_ai_container = None
        self.card_mode = False  # 是否在卡片复习模式

        # 初始化会话
        sessions = self.sm.get_all()
        if sessions:
            self.current_sid = sessions[0]["id"]
        else:
            self.current_sid = self.sm.create("默认课程")

        # 初始化视觉处理器
        vision_backend = self.config.get("vision_backend", "dashscope")
        vision_key = self.config.get("vision_api_key", "")
        if HAS_VISION:
            if vision_backend == "paddleocr":
                # PaddleOCR 不需要 API Key
                self.vision = DocumentVisionProcessor(backend="paddleocr")
            elif vision_key:
                self.vision = DocumentVisionProcessor(
                    backend=vision_backend, api_key=vision_key)
            else:
                self.vision = None
        else:
            self.vision = None

        self.kb = KnowledgeEngine(self.current_sid, vision_processor=self.vision)
        self.cards = StudyCardSystem(self.current_sid)

        # 窗口
        self.root = ttkb.Window(
            title="📚 DeepSeek 学习助手 v3.0",
            themename="litera",
            size=(1150, 740),
            minsize=(850, 520),
        )
        self._setup_ui()
        self._update_session_display()

        if self.config.get("api_key"):
            self.client = DeepSeekClient(
                self.config["api_key"],
                self.config.get("model", DEFAULT_MODEL)
            )
            self._update_status("已连接", "success")
        else:
            self._update_status("未连接", "secondary")

        self._refresh_kb_display()
        self.restore_conversation()
        self._add_system_message(
            "👋 欢迎使用 DeepSeek 学习助手 v3.0！\n"
            "📁 左侧管理课程和资料  |  ⚙ 右上角设置 API Key"
        )

        self.root.protocol("WM_DELETE_WINDOW", self._on_close)

    # ── 配置 ──────────────────────────────────────────────────────
    def _load_config(self) -> dict:
        if os.path.exists(CONFIG_FILE):
            try:
                with open(CONFIG_FILE, "r", encoding="utf-8") as f:
                    text = f.read()
                import yaml
                return yaml.safe_load(text) or {}
            except:
                pass
        return {"api_key": "", "model": DEFAULT_MODEL}

    def _save_config(self):
        import yaml
        with open(CONFIG_FILE, "w", encoding="utf-8") as f:
            yaml.dump(self.config, f, allow_unicode=True, default_flow_style=False)

    # ── UI 构建 ─────────────────────────────────────────────────
    def _setup_ui(self):
        self.root.grid_rowconfigure(0, weight=1)
        self.root.grid_columnconfigure(0, weight=1)

        main_pw = ttk.PanedWindow(self.root, orient=HORIZONTAL)
        main_pw.grid(row=0, column=0, sticky="nsew")

        # ═══ 左侧面板 ═══
        left = ttkb.Frame(main_pw, width=340, padding=0)
        left.pack_propagate(False)

        # 课程列表标题
        hdr = ttkb.Frame(left, padding=(12, 8))
        hdr.pack(fill=X)
        ttkb.Label(hdr, text="📚 课程列表",
                   font=("Microsoft YaHei UI", 12, "bold")).pack(side=LEFT)
        ttkb.Button(hdr, text="＋ 新课程", bootstyle="success-outline",
                    width=8, command=self._create_session).pack(side=RIGHT)

        # 会话列表
        slf = ttkb.Frame(left)
        slf.pack(fill=BOTH, expand=True, padx=8)
        self.s_canvas = tk.Canvas(slf, highlightthickness=0, bg="#f8f9fa")
        s_scroll = ttkb.Scrollbar(slf, orient=VERTICAL, command=self.s_canvas.yview)
        self.s_canvas.configure(yscrollcommand=s_scroll.set)
        self.s_inner = ttkb.Frame(self.s_canvas)
        self.s_inner.bind("<Configure>",
            lambda e: self.s_canvas.configure(scrollregion=self.s_canvas.bbox("all")))
        self.s_win = self.s_canvas.create_window(
            (0, 0), window=self.s_inner, anchor="nw")
        self.s_canvas.bind("<Configure>",
            lambda e: self.s_canvas.itemconfig(self.s_win, width=e.width))
        self.s_canvas.pack(side=LEFT, fill=BOTH, expand=True)
        s_scroll.pack(side=RIGHT, fill=Y)

        ttk.Separator(left, orient=HORIZONTAL).pack(fill=X, padx=12, pady=4)

        # ── 知识库区域 ──
        kbf = ttkb.Frame(left)
        kbf.pack(fill=X, padx=12, pady=(4, 0))
        ttkb.Label(kbf, text="📄 知识库",
                   font=("Microsoft YaHei UI", 10, "bold")).pack(side=LEFT)
        ttkb.Button(kbf, text="📤 上传", bootstyle="info-outline",
                    width=8, command=self._upload_file).pack(side=RIGHT)

        kbl = ttkb.Frame(left)
        kbl.pack(fill=BOTH, expand=True, padx=12, pady=4)
        self.k_canvas = tk.Canvas(kbl, height=150, highlightthickness=0, bg="#f8f9fa")
        k_scroll = ttkb.Scrollbar(kbl, orient=VERTICAL, command=self.k_canvas.yview)
        self.k_canvas.configure(yscrollcommand=k_scroll.set)
        self.k_inner = ttkb.Frame(self.k_canvas)
        self.k_inner.bind("<Configure>",
            lambda e: self.k_canvas.configure(scrollregion=self.k_canvas.bbox("all")))
        self.k_win = self.k_canvas.create_window(
            (0, 0), window=self.k_inner, anchor="nw")
        self.k_canvas.bind("<Configure>",
            lambda e: self.k_canvas.itemconfig(self.k_win, width=e.width))
        self.k_canvas.pack(side=LEFT, fill=BOTH, expand=True)
        k_scroll.pack(side=RIGHT, fill=Y)

        self.kb_stat = ttkb.Label(left, text="0 个文件",
            font=("Microsoft YaHei UI", 9), bootstyle="secondary")
        self.kb_stat.pack(fill=X, padx=12, pady=(0, 4))

        # ── 学习卡片区域 ──
        ttk.Separator(left, orient=HORIZONTAL).pack(fill=X, padx=12, pady=4)

        cf = ttkb.Frame(left)
        cf.pack(fill=X, padx=12, pady=(4, 4))
        ttkb.Label(cf, text="🎴 学习卡片",
                   font=("Microsoft YaHei UI", 10, "bold")).pack(side=LEFT)
        ttkb.Button(cf, text="＋新建", bootstyle="warning-outline",
                    width=6, command=self._add_card_dialog).pack(side=RIGHT)
        ttkb.Button(cf, text="📖复习", bootstyle="info-outline",
                    width=6, command=self._review_cards).pack(side=RIGHT, padx=3)

        self.card_list = ttkb.Frame(left)
        self.card_list.pack(fill=X, padx=12, pady=(0, 4))

        self.card_stat = ttkb.Label(left, text="0 张卡片",
            font=("Microsoft YaHei UI", 9), bootstyle="secondary")
        self.card_stat.pack(fill=X, padx=12, pady=(0, 8))

        main_pw.add(left, weight=0)

        # ═══ 右侧主区域 ═══
        right = ttkb.Frame(main_pw)
        right.grid_rowconfigure(2, weight=1)
        right.grid_columnconfigure(0, weight=1)

        # 标题栏
        tbar = ttkb.Frame(right, padding=(20, 10))
        tbar.grid(row=0, column=0, sticky="ew")
        tbar.grid_columnconfigure(0, weight=1)

        # 标题 + 学习统计
        title_left = ttkb.Frame(tbar)
        title_left.pack(side=LEFT)

        self.title_label = ttkb.Label(title_left, text="📖 默认课程",
            font=("Microsoft YaHei UI", 16, "bold"))
        self.title_label.pack(anchor="w")

        # 右侧按钮组
        btn_right = ttkb.Frame(tbar)
        btn_right.pack(side=RIGHT)

        self.stat_label = ttkb.Label(btn_right, text="● 未连接",
            font=("Microsoft YaHei UI", 9), bootstyle="secondary")
        self.stat_label.pack(side=LEFT, padx=(0, 10))

        # 主题切换按钮
        ttkb.Button(btn_right, text="🌙", bootstyle="outline-secondary",
                    width=3, command=self._toggle_theme).pack(side=RIGHT, padx=2)

        ttkb.Button(btn_right, text="⚙", bootstyle="outline-secondary",
                    width=3, command=self.open_settings).pack(side=RIGHT)

        # 分隔线
        ttk.Separator(right, orient=HORIZONTAL).grid(row=1, column=0, sticky="ew", padx=10)

        # 聊天区
        cf2 = ttkb.Frame(right, padding=(10, 5))
        cf2.grid(row=2, column=0, sticky="nsew")
        cf2.grid_rowconfigure(0, weight=1)
        cf2.grid_columnconfigure(0, weight=1)

        self.chat_canvas = tk.Canvas(cf2, highlightthickness=0, bg="#f0f2f5")
        c_scroll = ttk.Scrollbar(cf2, orient=VERTICAL, command=self.chat_canvas.yview)
        self.chat_canvas.configure(yscrollcommand=c_scroll.set)

        self.msg_frame = ttkb.Frame(self.chat_canvas, padding=10)
        self.msg_frame.bind("<Configure>",
            lambda e: self.chat_canvas.configure(scrollregion=self.chat_canvas.bbox("all")))

        self.c_win = self.chat_canvas.create_window(
            (0, 0), window=self.msg_frame, anchor="nw",
            width=self.chat_canvas.winfo_width())
        self.chat_canvas.bind("<Configure>",
            lambda e: self.chat_canvas.itemconfig(self.c_win, width=e.width))

        self.chat_canvas.grid(row=0, column=0, sticky="nsew")
        c_scroll.grid(row=0, column=1, sticky="ns")

        # 输入区
        inp = ttkb.Frame(right, padding=(10, 8))
        inp.grid(row=3, column=0, sticky="ew")
        inp.grid_columnconfigure(0, weight=1)

        self.input_text = scrolledtext.ScrolledText(
            inp, height=3, wrap=WORD,
            font=("Microsoft YaHei UI", 10),
            padx=10, pady=8, relief="flat", borderwidth=1)
        self.input_text.grid(row=0, column=0, sticky="ew", padx=(0, 8))
        self.input_text.bind("<Return>", self._on_enter)

        self.send_btn = ttkb.Button(inp, text="发送 ▶", bootstyle="success",
                                    command=self.send_message, width=10)
        self.send_btn.grid(row=0, column=1, sticky="ns")

        ttkb.Label(inp, text="Enter 发送 | Shift+Enter 换行",
            font=("Microsoft YaHei UI", 8), bootstyle="secondary"
        ).grid(row=1, column=0, columnspan=2, sticky="w", pady=(2, 0))

        main_pw.add(right, weight=1)

    # ═════════════════════════════════════════════════════════════
    #  会话管理
    # ═════════════════════════════════════════════════════════════

    def _update_session_display(self):
        for w in self.s_inner.winfo_children():
            w.destroy()

        sessions = self.sm.get_all()

        for sess in sessions:
            sid = sess["id"]
            active = sid == self.current_sid

            f = ttkb.Frame(self.s_inner, padding=(6, 5))
            f.pack(fill=X, pady=1)

            display_name = _truncate(sess["name"], 14)
            label = ttkb.Label(f,
                text=f"📁 {display_name}",
                font=("Microsoft YaHei UI", 10, "bold" if active else "normal"),
                bootstyle="inverse-primary" if active else "")
            label.pack(side=LEFT, fill=X, expand=True)

            # 消息计数小标签
            cnt = sess["message_count"]
            if cnt > 0:
                ttkb.Label(f, text=str(cnt),
                          font=("Microsoft YaHei UI", 8),
                          bootstyle="secondary").pack(side=RIGHT, padx=(2, 0))

            for w in [f, label]:
                w.bind("<Button-1>", lambda e, s=sid: self._switch_session(s))
                w.bind("<Button-3>", lambda e, s=sid, n=sess['name']: self._ctx_menu(e, s, n))

        self._update_title()

    def _ctx_menu(self, event, sid, name):
        menu = tk.Menu(self.root, tearoff=0, font=("Microsoft YaHei UI", 9))
        menu.add_command(label="✏️ 重命名", command=lambda: self._rename_session(sid))
        menu.add_command(label="🗑 删除", command=lambda: self._delete_session(sid))
        menu.tk_popup(event.x_root, event.y_root)

    def _create_session(self):
        sid = self.sm.create(f"课程 {len(self.sm.get_all()) + 1}")
        self.current_sid = sid
        self.kb = KnowledgeEngine(sid, vision_processor=self.vision)
        self.cards = StudyCardSystem(sid)
        self.conversation = []
        self._clear_chat_area()
        self._update_session_display()
        self._refresh_kb_display()
        self._refresh_cards()
        self._add_system_message(f"📖 已创建新课程，请上传学习资料")

    def _rename_session(self, sid):
        win = ttkb.Toplevel(self.root)
        win.title("重命名")
        win.geometry("350x120")
        win.transient(self.root)
        win.grab_set()

        ttkb.Label(win, text="新名称：", font=("Microsoft YaHei UI", 10)
                  ).pack(padx=20, pady=(15, 5))
        entry = ttkb.Entry(win, font=("Microsoft YaHei UI", 10))
        entry.pack(fill=X, padx=20)
        entry.insert(0, self.sm.get_name(sid))
        entry.select_range(0, END)
        entry.focus()

        def ok():
            n = entry.get().strip()
            if n:
                self.sm.rename(sid, n)
                self._update_session_display()
            win.destroy()

        ttkb.Button(win, text="确定", bootstyle="success",
                    command=ok).pack(pady=10)
        win.bind("<Return>", lambda e: ok())

    def _delete_session(self, sid):
        all_s = self.sm.get_all()
        if len(all_s) <= 1:
            messagebox.showwarning("提示", "至少保留一个课程")
            return
        if not messagebox.askyesno("确认", f"删除「{self.sm.get_name(sid)}」及其所有资料？"):
            return

        self.sm.delete(sid)
        remaining = self.sm.get_all()
        self.current_sid = remaining[0]["id"]
        self.kb = KnowledgeEngine(self.current_sid)
        self.cards = StudyCardSystem(self.current_sid)
        self.conversation = []
        self._clear_chat_area()
        self._update_session_display()
        self._refresh_kb_display()
        self._refresh_cards()
        self._add_system_message(f"已切换到「{self.sm.get_name(self.current_sid)}」")

    def _switch_session(self, sid):
        if sid == self.current_sid or self.is_loading:
            return
        self.current_sid = sid
        self.kb = KnowledgeEngine(sid, vision_processor=self.vision)
        self.cards = StudyCardSystem(sid)
        self.conversation = []
        self._clear_chat_area()
        self._update_session_display()
        self._refresh_kb_display()
        self._refresh_cards()
        self.restore_conversation()
        self._update_status(self._update_cost_display(), "success")

    def _update_title(self):
        self.title_label.configure(text=f"📖 {self.sm.get_name(self.current_sid)}")

    # ═════════════════════════════════════════════════════════════
    #  知识库
    # ═════════════════════════════════════════════════════════════

    def _refresh_kb_display(self):
        for w in self.k_inner.winfo_children():
            w.destroy()

        files = self.kb.get_files()

        if not files:
            ttkb.Label(self.k_inner, text="📭 暂无资料\n点击「上传」添加",
                      font=("Microsoft YaHei UI", 9), bootstyle="secondary",
                      justify=CENTER).pack(pady=15)
        else:
            for f in files:
                ff = ttkb.Frame(self.k_inner)
                ff.pack(fill=X, pady=1)

                ext = os.path.splitext(f["file_name"])[1].lower()
                icon = FILE_ICONS.get(ext, "📄")

                # 截断长文件名
                display_name = _truncate(f["file_name"], 16)

                ttkb.Label(ff, text=f"{icon} {display_name}",
                          font=("Microsoft YaHei UI", 9)
                ).pack(side=LEFT, fill=X, expand=True)

                ttkb.Label(ff, text=f"{f['chunk_count']}段",
                          font=("Microsoft YaHei UI", 8),
                          bootstyle="secondary"
                ).pack(side=RIGHT, padx=(2, 0))

                ttkb.Button(ff, text="✕", bootstyle="danger-outline",
                           width=2, padding=(2, 0),
                           command=lambda n=f['file_name']: self._remove_file(n)
                ).pack(side=RIGHT)

        fc, cc, tc = self.kb.get_stats()
        size_hint = f"（{tc//1000}k字）" if tc > 1000 else ""
        self.kb_stat.configure(text=f"{fc} 个文件 · {cc} 个片段{size_hint}")

    def _upload_file(self):
        types = [
            ("支持的文件", "*.txt *.md *.pdf *.docx *.pptx"),
            ("文本文件", "*.txt"), ("Markdown", "*.md"),
            ("PDF", "*.pdf"), ("Word", "*.docx"),
            ("PowerPoint", "*.pptx"), ("所有文件", "*.*"),
        ]
        files = filedialog.askopenfilenames(title="选择学习资料", filetypes=types)
        if not files:
            return

        for f in files:
            ok, msg = self.kb.add_file(f)
            if not ok:
                messagebox.showwarning("导入失败", msg)

        self._refresh_kb_display()
        fn = os.path.basename(files[0]) if len(files) == 1 else f"{len(files)} 个文件"
        self._add_system_message(f"📎 已导入 {fn}，可以提问了")

    def _remove_file(self, name):
        if messagebox.askyesno("确认", f"从知识库移除「{name}」？"):
            self.kb.remove_file(name)
            self._refresh_kb_display()
            self._add_system_message(f"🗑 已移除「{name}」")

    # ═════════════════════════════════════════════════════════════
    #  学习卡片（借鉴 Hermes Skills 系统）
    # ═════════════════════════════════════════════════════════════

    def _refresh_cards(self):
        for w in self.card_list.winfo_children():
            w.destroy()

        cards = self.cards.get_cards()
        if not cards:
            ttkb.Label(self.card_list, text="暂无卡片，点「＋新建」创建",
                      font=("Microsoft YaHei UI", 9), bootstyle="secondary"
            ).pack(pady=5)
        else:
            for c in cards[:5]:  # 只显示最近5张
                cf = ttkb.Frame(self.card_list)
                cf.pack(fill=X, pady=1)

                icon = CARD_ICONS.get(c["category"], "📌")
                title = c["title"][:20] + ("…" if len(c["title"]) > 20 else "")
                ttkb.Label(cf,
                    text=f"{icon} {title}",
                    font=("Microsoft YaHei UI", 9)
                ).pack(side=LEFT)

                ttkb.Label(cf,
                    text=f"✕{c['review_count']}",
                    font=("Microsoft YaHei UI", 8),
                    bootstyle="secondary"
                ).pack(side=RIGHT)

        cnt = len(cards)
        cats = self.cards.get_categories()
        self.card_stat.configure(text=f"{cnt} 张卡片 · {len(cats)} 分类")

    def _add_card_dialog(self):
        win = ttkb.Toplevel(self.root)
        win.title("新建学习卡片")
        win.geometry("500x400")
        win.transient(self.root)
        win.grab_set()

        f = ttkb.Frame(win, padding=15)
        f.pack(fill=BOTH, expand=True)

        ttkb.Label(f, text="🎴 新建学习卡片",
                   font=("Microsoft YaHei UI", 12, "bold")).pack(anchor="w", pady=(0, 10))

        ttkb.Label(f, text="标题 *", font=("Microsoft YaHei UI", 10)).pack(anchor="w")
        title_e = ttkb.Entry(f, font=("Microsoft YaHei UI", 10))
        title_e.pack(fill=X, pady=(2, 8))

        ttkb.Label(f, text="问题（可选）", font=("Microsoft YaHei UI", 10)).pack(anchor="w")
        q_e = scrolledtext.ScrolledText(f, height=3, font=("Microsoft YaHei UI", 10))
        q_e.pack(fill=X, pady=(2, 8))

        ttkb.Label(f, text="答案/知识点 *", font=("Microsoft YaHei UI", 10)).pack(anchor="w")
        a_e = scrolledtext.ScrolledText(f, height=4, font=("Microsoft YaHei UI", 10))
        a_e.pack(fill=X, pady=(2, 8))

        # 分类
        bf = ttkb.Frame(f)
        bf.pack(fill=X, pady=5)
        ttkb.Label(bf, text="分类：", font=("Microsoft YaHei UI", 10)).pack(side=LEFT)
        cat_e = ttkb.Combobox(bf,
            values=["通用", "知识点", "公式", "定义", "代码", "错题"],
            state="normal", font=("Microsoft YaHei UI", 10), width=12)
        cat_e.pack(side=LEFT, padx=5)
        cat_e.set("通用")

        ttkb.Label(bf, text="标签：", font=("Microsoft YaHei UI", 10)).pack(side=LEFT, padx=(10, 0))
        tag_e = ttkb.Entry(bf, font=("Microsoft YaHei UI", 10), width=15)
        tag_e.pack(side=LEFT, padx=5)

        def save():
            title = title_e.get().strip()
            answer = a_e.get("1.0", "end-1c").strip()
            if not title or not answer:
                messagebox.showwarning("提示", "标题和答案为必填")
                return
            question = q_e.get("1.0", "end-1c").strip()
            self.cards.add_card(title, answer, question, cat_e.get(), tag_e.get())
            self._refresh_cards()
            win.destroy()
            self._add_system_message(f"🎴 已创建卡片「{title}」")

        ttkb.Button(f, text="💾 保存卡片", bootstyle="success",
                    command=save).pack(pady=10)

    def _review_cards(self):
        cards = self.cards.get_cards()
        if not cards:
            messagebox.showinfo("提示", "还没有学习卡片，先创建一些吧！")
            return

        self._clear_chat_area()
        self._add_system_message(f"🎴 开始复习！共 {len(cards)} 张卡片，逐个展示中")

        for c in cards:
            self._show_card(c)

    def _show_card(self, card: dict):
        """展示一张卡片"""
        icon = CARD_ICONS.get(card["category"], "📌")

        # 卡片标题
        cf = ttkb.Frame(self.msg_frame)
        cf.pack(fill=X, pady=5, padx=10)

        inner = ttkb.Frame(cf)
        inner.pack(anchor="w")

        card_frame = ttkb.Frame(inner, bootstyle="warning", padding=8)
        card_frame.pack()

        ttkb.Label(card_frame, text=f"{icon} {card['title']}  [{card['category']}]",
                  font=("Microsoft YaHei UI", 11, "bold")).pack(anchor="w")

        if card.get("question"):
            ttkb.Label(card_frame, text=f"❓ {card['question']}",
                      font=("Microsoft YaHei UI", 10),
                      wraplength=400).pack(anchor="w", pady=(5, 0))

        # 答案（用 Text 组件显示多行）
        ans = tk.Text(card_frame, wrap=WORD,
                     font=("Microsoft YaHei UI", 10),
                     bg="#fff8e1", relief="flat",
                     padx=8, pady=5, height=3,
                     highlightthickness=0)
        ans.pack(fill=X, pady=(5, 0))
        ans.insert("1.0", f"💡 {card['answer']}")
        ans.configure(state=DISABLED)

        # 复习按钮
        btn_f = ttkb.Frame(card_frame)
        btn_f.pack(fill=X, pady=(5, 0))

        def rate(d):
            self.cards.review_card(card["id"], d)
            self._refresh_cards()
            self._scroll_to_bottom()

        ttkb.Button(btn_f, text="😊 简单", bootstyle="success-outline",
                    width=8, command=lambda: rate(1)).pack(side=LEFT, padx=2)
        ttkb.Button(btn_f, text="🤔 中等", bootstyle="warning-outline",
                    width=8, command=lambda: rate(3)).pack(side=LEFT, padx=2)
        ttkb.Button(btn_f, text="😰 困难", bootstyle="danger-outline",
                    width=8, command=lambda: rate(5)).pack(side=LEFT, padx=2)

        self._scroll_to_bottom()

    # ═════════════════════════════════════════════════════════════
    #  聊天
    # ═════════════════════════════════════════════════════════════

    def _clear_chat_area(self):
        for w in self.msg_frame.winfo_children():
            w.destroy()

    def _add_system_message(self, text):
        f = ttkb.Frame(self.msg_frame)
        f.pack(fill=X, pady=4, padx=5)
        inner = ttkb.Frame(f)
        inner.pack(anchor="center")
        ttkb.Label(inner, text=f"📌 {text}",
                  font=("Microsoft YaHei UI", 9), bootstyle="secondary"
        ).pack(padx=10, pady=5)

    def _create_bubble(self, role, content, stream=False):
        """创建聊天气泡 — 现代圆角设计"""
        user = role == "user"
        f = ttkb.Frame(self.msg_frame)
        f.pack(fill=X, pady=(0, 10), padx=8)
        inner = ttkb.Frame(f)
        inner.pack(anchor="e" if user else "w")

        tag = "👤 你" if user else "🤖 DeepSeek"
        tag_color = "#2e7d32" if user else "#1565c0"

        # 角色标签
        ttkb.Label(inner, text=tag,
                  font=("Microsoft YaHei UI", 9, "bold"),
                  foreground=tag_color
        ).pack(anchor="w" if not user else "e", padx=6, pady=(0, 4))

        # 气泡主体 — 使用 Frame 模拟圆角
        bg_color = "#e8f5e9" if user else "#ffffff"
        border_color = "#c8e6c9" if user else "#e0e0e0"

        bubble = tk.Frame(inner, bg=bg_color,
                         highlightbackground=border_color,
                         highlightthickness=1)
        bubble.pack()

        tw = tk.Text(bubble, wrap=WORD, font=("Microsoft YaHei UI", 10),
                    bg=bg_color, relief="flat", padx=16, pady=12,
                    height=3, width=65, highlightthickness=0,
                    borderwidth=0)
        tw.pack()
        tw.insert("1.0", content)
        if not stream:
            tw.configure(state=DISABLED)

        # 时间戳
        ts = datetime.now().strftime("%H:%M")
        ttkb.Label(inner, text=ts, font=("Microsoft YaHei UI", 8),
                  bootstyle="secondary"
        ).pack(anchor="e" if user else "w", padx=6, pady=(3, 0))

        return tw

    def send_message(self):
        if self.is_loading:
            return

        text = self.input_text.get("1.0", "end-1c").strip()
        if not text:
            return

        if not self.client or not self.config.get("api_key"):
            messagebox.showwarning("提示", "请先设置 API Key")
            self.open_settings()
            return

        self.input_text.delete("1.0", END)

        self._create_bubble("user", text)
        self.conversation.append({"role": "user", "content": text})
        self.save_message("user", text)

        self.root.after(50, self._scroll_to_bottom)

        self.is_loading = True
        self.send_btn.configure(text="响应中...", state=DISABLED)
        self._update_status("思考中...", "warning")

        threading.Thread(target=self._do_chat, args=(text,), daemon=True).start()

    def _do_chat(self, user_text):
        try:
            # 1. 检索知识库（RAG）
            context = self.kb.format_context(user_text)

            # 2. 构建三层系统提示
            fc, cc, tc = self.kb.get_stats()
            card_count = len(self.cards.get_cards())
            q_count = len(self.conversation)

            if context and "未找到相关内容" not in context:
                context_prompt = KNOWLEDGE_CONTEXT_TEMPLATE.format(context=context)
            else:
                context_prompt = "\n【参考资料】\n（当前课程暂无相关学习资料）"

            volatile = VOLATILE_NOTE.format(
                timestamp=datetime.now().strftime("%Y-%m-%d %H:%M"),
                q_count=q_count,
                file_count=fc,
                card_count=card_count,
            )

            sys_prompt = STABLE_IDENTITY + context_prompt + volatile

            # 3. 构建消息列表
            messages = [{"role": "system", "content": sys_prompt}]
            messages.extend(self.conversation)

            # 4. 流式输出（支持 reasoning 显示）
            use_reasoner = "reasoner" in self.config.get("model", "")
            self.root.after(0, self._create_stream_bubble)

            full = ""
            reasoning_text = ""

            for chunk_type, chunk in self.client.chat_stream(messages):
                if chunk_type == "reasoning":
                    reasoning_text += chunk
                elif chunk_type == "content":
                    if reasoning_text and not full:
                        # 第一次收到 content，表示推理结束，可以展示了
                        pass
                    full += chunk
                    self.root.after(0, self._update_stream, full)

            # 如果只有推理没有内容（极端情况），把推理当内容
            if not full and reasoning_text:
                full = f"[推理过程]\n{reasoning_text}"

            self.root.after(0, self._finish_chat, full)
            self.sm.update_time(self.current_sid)
            self._update_session_display()

        except requests.exceptions.HTTPError as e:
            msg = f"API 请求失败 ({e.response.status_code})"
            if e.response.status_code == 401:
                msg = "API Key 无效，请检查设置"
            self.root.after(0, self._handle_error, msg)
        except requests.exceptions.ConnectionError:
            self.root.after(0, self._handle_error, "无法连接 DeepSeek API，请检查网络")
        except Exception as e:
            self.root.after(0, self._handle_error, f"错误：{e}")

    def _create_stream_bubble(self):
        f = ttkb.Frame(self.msg_frame)
        f.pack(fill=X, pady=(0, 10), padx=8)
        c = ttkb.Frame(f)
        c.pack(anchor="w")

        ttkb.Label(c, text="🤖 DeepSeek",
                  font=("Microsoft YaHei UI", 9, "bold"),
                  foreground="#1565c0").pack(anchor="w", padx=6, pady=(0, 4))

        bubble = tk.Frame(c, bg="#ffffff",
                         highlightbackground="#e0e0e0",
                         highlightthickness=1)
        bubble.pack()

        tw = tk.Text(bubble, wrap=WORD, font=("Microsoft YaHei UI", 10),
                    bg="#ffffff", relief="flat", padx=16, pady=12,
                    height=3, width=65, highlightthickness=0,
                    borderwidth=0)
        tw.pack()
        tw.configure(state="normal")
        tw.insert("1.0", "▊")

        self.current_ai_widget = tw
        self.current_ai_container = c

        self.root.after(300, self._blink)
        self._scroll_to_bottom()

    def _blink(self):
        if self.current_ai_widget and self.is_loading:
            c = self.current_ai_widget.get("1.0", "end-1c")
            if c.endswith("▊"):
                self.current_ai_widget.delete("end-2c", "end-1c")
                self.current_ai_widget.insert(END, "▌")
            elif c.endswith("▌"):
                self.current_ai_widget.delete("end-2c", "end-1c")
                self.current_ai_widget.insert(END, "▊")
            self.root.after(400, self._blink)

    def _update_stream(self, text):
        if self.current_ai_widget:
            self.current_ai_widget.configure(state="normal")
            self.current_ai_widget.delete("1.0", END)
            self.current_ai_widget.insert("1.0", text + "▊")
            lines = int(self.current_ai_widget.index("end-1c").split(".")[0])
            self.current_ai_widget.configure(height=min(max(lines + 1, 3), 22))
            self._scroll_to_bottom()

    def _finish_chat(self, text):
        if self.current_ai_widget:
            self.current_ai_widget.configure(state="normal")
            self.current_ai_widget.delete("1.0", END)
            self.current_ai_widget.insert("1.0", text)
            self.current_ai_widget.configure(state=DISABLED)

            ttkb.Label(self.current_ai_container, text=datetime.now().strftime("%H:%M"),
                      font=("Microsoft YaHei UI", 8), bootstyle="secondary"
            ).pack(anchor="w", padx=6, pady=(3, 0))

        self.conversation.append({"role": "assistant", "content": text})
        self.save_message("assistant", text)
        self._update_status(self._update_cost_display(), "success")
        self._reset_state()

    def _handle_error(self, msg):
        if self.current_ai_widget:
            self.current_ai_widget.configure(state="normal")
            self.current_ai_widget.delete("1.0", END)
            self.current_ai_widget.insert("1.0", f"❌ {msg}")
            self.current_ai_widget.configure(state=DISABLED)

            ttkb.Label(self.current_ai_container, text=datetime.now().strftime("%H:%M"),
                      font=("Microsoft YaHei UI", 8), bootstyle="secondary"
            ).pack(anchor="w", padx=6, pady=(3, 0))

        self._reset_state()

    def _reset_state(self):
        self.is_loading = False
        self.send_btn.configure(text="发送 ▶", state=NORMAL)
        self.current_ai_widget = None
        self.current_ai_container = None
        self._update_status("就绪", "success")

    def _update_status(self, text, style="secondary"):
        self.stat_label.configure(text=f"● {text}", bootstyle=style)

    def _scroll_to_bottom(self):
        self.chat_canvas.yview_moveto(1.0)
        self.root.update_idletasks()

    def _on_enter(self, event):
        if event.state & 0x1:
            return
        self.send_message()
        return "break"

    # ═════════════════════════════════════════════════════════════
    #  设置
    # ═════════════════════════════════════════════════════════════

    def open_settings(self):
        win = ttkb.Toplevel(self.root)
        win.title("设置")
        win.geometry("480x460")
        win.resizable(False, False)
        win.transient(self.root)
        win.grab_set()

        f = ttkb.Frame(win, padding=20)
        f.pack(fill=BOTH, expand=True)

        ttkb.Label(f, text="⚙ 设置", font=("Microsoft YaHei UI", 14, "bold")
                  ).pack(anchor="w", pady=(0, 15))

        # DeepSeek
        ttkb.Label(f, text="🤖 DeepSeek API Key", font=("Microsoft YaHei UI", 10)
                  ).pack(anchor="w")
        entry = ttkb.Entry(f, font=("Consolas", 10), show="*")
        entry.pack(fill=X, pady=(4, 3))
        entry.insert(0, self.config.get("api_key", ""))
        ttkb.Label(f, text="platform.deepseek.com 获取",
                  font=("Microsoft YaHei UI", 8), bootstyle="info"
        ).pack(anchor="w", pady=(0, 8))

        ttkb.Label(f, text="模型选择", font=("Microsoft YaHei UI", 10)).pack(anchor="w")
        combo = ttkb.Combobox(f,
            values=["deepseek-v4-flash ⚡（快速省钱）", "deepseek-v4-pro 👑（最强旗舰）"],
            state="readonly", font=("Microsoft YaHei UI", 10))
        combo.pack(fill=X, pady=(4, 3))
        # 映射配置值到显示文本
        model_display_map = {
            "deepseek-v4-flash": "deepseek-v4-flash ⚡（快速省钱）",
            "deepseek-v4-pro": "deepseek-v4-pro 👑（最强旗舰）",
        }
        current_model = self.config.get("model", DEFAULT_MODEL)
        combo.set(model_display_map.get(current_model, "deepseek-v4-flash ⚡（快速省钱）"))
        ttkb.Label(f, text="v4-flash 适合日常学习 | v4-pro 适合复杂推理",
                  font=("Microsoft YaHei UI", 8), bootstyle="secondary"
        ).pack(anchor="w", pady=(0, 10))

        # 图片理解（国内用通义千问 VL，国外用 Gemini）
        ttk.Separator(f, orient=HORIZONTAL).pack(fill=X, pady=5)

        ttkb.Label(f, text="🖼️ 图片理解（可选）— 用于翻译文档中的图片/图表",
                  font=("Microsoft YaHei UI", 10)
        ).pack(anchor="w")

        # 后端选择
        backend_f = ttkb.Frame(f)
        backend_f.pack(fill=X, pady=(5, 3))
        ttkb.Label(backend_f, text="后端：", font=("Microsoft YaHei UI", 10)).pack(side=LEFT)
        vision_combo = ttkb.Combobox(backend_f,
            values=["dashscope（通义千问 - 国内推荐 ⭐）",
                    "gemini（Google - 海外可用）",
                    "paddleocr（本地方案）"],
            state="readonly", font=("Microsoft YaHei UI", 10), width=32)
        vision_combo.pack(side=LEFT, padx=5)

        # 设置当前后端
        current_backend = self.config.get("vision_backend", "dashscope")
        if current_backend == "dashscope":
            vision_combo.set("dashscope（通义千问 - 国内推荐 ⭐）")
        elif current_backend == "gemini":
            vision_combo.set("gemini（Google - 海外可用）")
        elif current_backend == "paddleocr":
            vision_combo.set("paddleocr（本地方案）")
        else:
            vision_combo.set("dashscope（通义千问 - 国内推荐 ⭐）")

        # API Key 输入
        vision_entry = ttkb.Entry(f, font=("Consolas", 10), show="*")
        vision_entry.pack(fill=X, pady=(4, 3))
        vision_entry.insert(0, self.config.get("vision_api_key", ""))

        # 提示文字（动态更新）
        hint_label = ttkb.Label(f,
            text="dashscope.aliyun.com 获取 | 国内推荐，理解图表+文字",
            font=("Microsoft YaHei UI", 8), bootstyle="info")
        hint_label.pack(anchor="w", pady=(0, 5))

        def update_hint(*args):
            sel = vision_combo.get()
            if "dashscope" in sel:
                hint_label.configure(text="dashscope.aliyun.com 获取 | 国内推荐 ⭐ 理解图表+文字")
            elif "gemini" in sel:
                hint_label.configure(text="aistudio.google.com 获取 | 理解图表+文字（国内不可用）")
            elif "paddleocr" in sel:
                hint_label.configure(text="pip install paddleocr | 本地方案，仅识别文字，不需要 API Key")
        vision_combo.bind("<<ComboboxSelected>>", update_hint)
        update_hint()

        bf = ttkb.Frame(f)
        bf.pack(fill=X)
        ttkb.Button(bf, text="💾 保存", bootstyle="success",
                   command=lambda: self._save_settings(
                       win, entry.get().strip(), combo.get(),
                       vision_combo.get(), vision_entry.get().strip()),
                   width=12).pack(side=RIGHT, padx=(5, 0))
        ttkb.Button(bf, text="取消", bootstyle="secondary-outline",
                   command=win.destroy, width=12).pack(side=RIGHT)

    def _save_settings(self, win, api_key, model_display, vision_backend_str, vision_key=""):
        if not api_key:
            messagebox.showwarning("提示", "请输入 DeepSeek API Key")
            return

        # 解析模型名（从显示文本到配置值）
        model_parse_map = {
            "flash": "deepseek-v4-flash",
            "pro": "deepseek-v4-pro",
        }
        model = "deepseek-v4-flash"
        for key, val in model_parse_map.items():
            if key in model_display.lower():
                model = val
                break

        self.config["api_key"] = api_key
        self.config["model"] = model

        # 解析视觉后端
        backend_map = {
            "dashscope": "dashscope",
            "gemini": "gemini",
            "paddleocr": "paddleocr",
        }
        chosen_backend = "dashscope"
        for key, val in backend_map.items():
            if key in vision_backend_str.lower():
                chosen_backend = val
                break
        self.config["vision_backend"] = chosen_backend
        self.config["vision_api_key"] = vision_key

        self._save_config()
        self.client = DeepSeekClient(api_key, model)

        # 初始化视觉处理器
        if chosen_backend == "paddleocr":
            if HAS_VISION:
                self.vision = DocumentVisionProcessor(backend="paddleocr")
                self.kb.vision = self.vision
            else:
                self.vision = None
                self.kb.vision = None
        elif vision_key and HAS_VISION:
            self.vision = DocumentVisionProcessor(
                backend=chosen_backend, api_key=vision_key)
            self.kb.vision = self.vision
        else:
            self.vision = None
            if self.kb:
                self.kb.vision = None
        self._update_status("已连接", "success")
        win.destroy()

    # ════════════════════════════════════════════════════════════════
    #  对话持久化 — 保存/加载到 SQLite
    # ════════════════════════════════════════════════════════════════

    def save_message(self, role: str, content: str):
        """保存一条消息到数据库"""
        db = Database()
        db.execute(
            "INSERT INTO messages (session_id, role, content, created_at) VALUES (?, ?, ?, ?)",
            (self.current_sid, role, content, datetime.now().isoformat())
        )
        db.commit()
        # 粗略 token 计数 (中文约1.5字/token, 英文约4字/token)
        chars = len(content)
        tokens = max(1, chars // 2)
        return tokens

    def load_messages(self, limit: int = 50) -> list[dict]:
        """从数据库加载当前会话的最近消息"""
        db = Database()
        rows = db.fetchall(
            "SELECT role, content, created_at FROM messages WHERE session_id=? ORDER BY id ASC LIMIT ?",
            (self.current_sid, limit)
        )
        return [{"role": r["role"], "content": r["content"]} for r in rows]

    def restore_conversation(self):
        """恢复当前会话的历史对话到聊天界面"""
        messages = self.load_messages()
        if not messages:
            return

        for msg in messages:
            if msg["role"] == "user":
                self._create_bubble("user", msg["content"])
            elif msg["role"] == "assistant":
                self._create_bubble("assistant", msg["content"])
        self._scroll_to_bottom()

    # ════════════════════════════════════════════════════════════════
    #  Token / 费用跟踪
    # ════════════════════════════════════════════════════════════════

    def _approx_tokens(self, text: str) -> int:
        """粗略估算 token 数"""
        if not text:
            return 0
        # 中文字符约 1.5 字/token，英文约 4 字/token
        chinese = len(re.findall(r'[\u4e00-\u9fff]', text))
        other = len(text) - chinese
        return max(1, int(chinese * 1.5 + other / 4))

    def _update_cost_display(self):
        """更新费用显示"""
        # 从数据库统计当前会话的 token 使用
        db = Database()
        rows = db.fetchall(
            "SELECT role, SUM(LENGTH(content)) as chars FROM messages WHERE session_id=? GROUP BY role",
            (self.current_sid,)
        )
        input_tokens = 0
        output_tokens = 0
        for r in rows:
            tokens = self._approx_tokens(r["chars"] or "")
            if r["role"] == "assistant":
                output_tokens += tokens
            else:
                input_tokens += tokens

        # DeepSeek v4-flash 约 0.5元/百万输入token, 2元/百万输出token
        cost = (input_tokens / 1_000_000 * 0.5 + output_tokens / 1_000_000 * 2)
        total_tokens = input_tokens + output_tokens

        status_text = f"对话 {len(self.conversation)//2} 轮"
        if total_tokens > 0:
            status_text += f" | {total_tokens/1000:.1f}K tokens"
        if cost > 0.01:
            status_text += f" | ≈¥{cost:.2f}"

        return status_text

    def _on_close(self):
        self.root.destroy()

    # ════════════════════════════════════════════════════════════════
    #  主题切换
    # ════════════════════════════════════════════════════════════════

    def _toggle_theme(self):
        """切换深色/浅色主题"""
        current = self.root.style.theme_use()
        if current in ("litera", "flatly", "cosmo", "journal", "sandstone"):
            self.root.style.theme_use("darkly")
            # 更新主题按钮
            for w in self.root.winfo_children():
                self._update_theme_recursive(w)
        else:
            self.root.style.theme_use("litera")

    def _update_theme_recursive(self, widget):
        """递归更新组件主题（Canvas 背景色）"""
        if isinstance(widget, tk.Canvas):
            current = self.root.style.theme_use()
            if current == "darkly":
                widget.configure(bg="#2b2b2b")
            else:
                widget.configure(bg="#f0f2f5" if widget == self.chat_canvas else "#f8f9fa")

    def run(self):
        self.root.mainloop()


# ══════════════════════════════════════════════════════════════════════
#  入口
# ══════════════════════════════════════════════════════════════════════

if __name__ == "__main__":
    # 检查依赖
    deps_ok = True
    for mod in ["yaml"]:
        try:
            __import__(mod)
        except ImportError:
            print(f"❌ 缺少依赖: {mod}，请运行: pip install pyyaml")
            deps_ok = False

    if deps_ok:
        app = DeepSeekLearnerApp()
        app.run()
